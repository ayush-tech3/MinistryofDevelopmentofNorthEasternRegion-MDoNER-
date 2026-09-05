"""
SIH 2026 Official 6-Slide Presentation Generator for Team AlertNex
Theme: Disaster Management | Category: Software
Problem Statement ID: SIH26001
Organization: Ministry of Development of North Eastern Region (MDoNER)
Background: Clean White / Light Grey SIH Professional Style
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Professional Clean SIH Color Palette ──
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # Pure White
BG_LIGHT_GREY = RGBColor(0xF8, 0xFA, 0xFC)     # Subtle Light Grey Card Fill
CARD_BORDER = RGBColor(0xCB, 0xD5, 0xE1)       # Crisp Light Slate Border
DARK_NAVY = RGBColor(0x0A, 0x19, 0x2F)         # Primary Dark Navy
SLATE_BLUE = RGBColor(0x1E, 0x3A, 0x5F)        # Secondary Deep Blue
SIH_ORANGE = RGBColor(0xFF, 0x6B, 0x00)        # Official SIH Orange
FOREST_GREEN = RGBColor(0x00, 0x8A, 0x4B)      # Clean Green Accent
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)         # Charcoal Dark Text
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)        # Slate Grey Subtext
HEADER_LINE = RGBColor(0xE2, 0xE8, 0xF0)       # Divider Line

# Risk Category Colors
RISK_GREEN = RGBColor(0x10, 0xB9, 0x81)        # Low Risk
RISK_YELLOW = RGBColor(0xF5, 0x9E, 0x0B)       # Moderate Risk
RISK_ORANGE = RGBColor(0xF9, 0x73, 0x16)       # High Risk
RISK_RED = RGBColor(0xEF, 0x44, 0x44)          # Critical Risk

# Widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

def apply_clean_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WHITE
    bg.line.fill.background()
    return bg

def add_slide_header_footer(slide, slide_num, title_text):
    # Top Header Left: SMART INDIA HACKATHON 2026
    top_left = slide.shapes.add_textbox(Inches(0.8), Inches(0.28), Inches(4.5), Inches(0.35))
    tf_tl = top_left.text_frame
    p_tl = tf_tl.paragraphs[0]
    p_tl.text = "SMART INDIA HACKATHON 2026"
    p_tl.font.bold = True
    p_tl.font.size = Pt(11)
    p_tl.font.color.rgb = SIH_ORANGE

    # Top Header Right: Professional SIH-style branding area
    top_right = slide.shapes.add_textbox(Inches(7.2), Inches(0.28), Inches(5.33), Inches(0.35))
    tf_tr = top_right.text_frame
    p_tr = tf_tr.paragraphs[0]
    p_tr.text = "PS ID: SIH26001  |  MDoNER  |  Software Edition"
    p_tr.font.bold = True
    p_tr.font.size = Pt(10)
    p_tr.font.color.rgb = DARK_NAVY
    p_tr.alignment = PP_ALIGN.RIGHT

    # Thin Divider Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.68), Inches(11.733), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = HEADER_LINE
    line.line.fill.background()

    # Main Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.78), Inches(11.733), Inches(0.55))
    tf_t = title_box.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.font.bold = True
    p_t.font.size = Pt(22)
    p_t.font.color.rgb = DARK_NAVY

    # Bottom Footer Left
    footer_left = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(9.5), Inches(0.32))
    tf_fl = footer_left.text_frame
    p_fl = tf_fl.paragraphs[0]
    p_fl.text = "Team AlertNex | Smart India Hackathon 2026 | PS: SIH26001 | Disaster Management"
    p_fl.font.size = Pt(9.5)
    p_fl.font.color.rgb = TEXT_MUTED

    # Bottom Footer Right: Slide X of 6
    footer_right = slide.shapes.add_textbox(Inches(10.5), Inches(7.05), Inches(2.033), Inches(0.32))
    tf_fr = footer_right.text_frame
    p_fr = tf_fr.paragraphs[0]
    p_fr.text = f"Slide {slide_num} of 6"
    p_fr.font.bold = True
    p_fr.font.size = Pt(9.5)
    p_fr.font.color.rgb = DARK_NAVY
    p_fr.alignment = PP_ALIGN.RIGHT

# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide1 = prs.slides.add_slide(blank_layout)
apply_clean_background(slide1)

# Top Bar Branding
tb = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.4))
tb.fill.solid()
tb.fill.fore_color.rgb = BG_LIGHT_GREY
tb.line.color.rgb = CARD_BORDER
tb.line.width = Pt(1)
tb_tf = tb.text_frame
tb_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tb_p = tb_tf.paragraphs[0]
tb_p.text = "SMART INDIA HACKATHON 2026  •  COLLEGE INTERNAL EVALUATION ROUND"
tb_p.font.bold = True
tb_p.font.size = Pt(11)
tb_p.font.color.rgb = SIH_ORANGE
tb_p.alignment = PP_ALIGN.CENTER

# Main Title & Subtitle Box
m_box = slide1.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.733), Inches(1.2))
m_tf = m_box.text_frame
m_tf.word_wrap = True
p = m_tf.paragraphs[0]
p.text = "AI-Based Early Warning & Landslide Monitoring System"
p.font.bold = True
p.font.size = Pt(27)
p.font.color.rgb = DARK_NAVY

p_sub = m_tf.add_paragraph()
p_sub.text = "for North Eastern Region (NER) of India"
p_sub.font.bold = True
p_sub.font.size = Pt(17)
p_sub.font.color.rgb = SIH_ORANGE

# Project Details Pill Cards (Left Column)
det_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.35), Inches(4.3), Inches(4.55))
det_card.fill.solid()
det_card.fill.fore_color.rgb = BG_LIGHT_GREY
det_card.line.color.rgb = CARD_BORDER
det_card.line.width = Pt(1.2)
d_tf = det_card.text_frame
d_tf.word_wrap = True
d_tf.margin_left = Inches(0.22)
d_tf.margin_top = Inches(0.2)

p = d_tf.paragraphs[0]
p.text = "PROJECT SPECIFICATIONS"
p.font.bold = True
p.font.size = Pt(12)
p.font.color.rgb = DARK_NAVY

details = [
    ("Problem Statement ID", "SIH26001"),
    ("Theme", "Disaster Management"),
    ("Category", "Software Edition"),
    ("Ministry / Organization", "Ministry of Development of North Eastern Region (MDoNER)"),
    ("Team Name", "AlertNex"),
    ("Team Leader", "Ayush Kumar"),
    ("Core Focus", "Dynamic Risk + Connectivity Impact + Offline Ground Sync")
]
for lbl, val in details:
    p = d_tf.add_paragraph()
    p.text = f"{lbl}:"
    p.font.bold = True
    p.font.size = Pt(9.8)
    p.font.color.rgb = DARK_NAVY
    p_val = d_tf.add_paragraph()
    p_val.text = f"  {val}"
    p_val.font.size = Pt(9.5)
    p_val.font.color.rgb = TEXT_DARK

# Team Members Table Card (Right Column)
team_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(2.35), Inches(7.233), Inches(4.55))
team_card.fill.solid()
team_card.fill.fore_color.rgb = BG_LIGHT_GREY
team_card.line.color.rgb = CARD_BORDER
team_card.line.width = Pt(1.2)
t_tf = team_card.text_frame
t_tf.word_wrap = True
t_tf.margin_left = Inches(0.25)
t_tf.margin_top = Inches(0.18)

p = t_tf.paragraphs[0]
p.text = "TEAM ALERTSPEC PROFILE — ALL 6 MEMBERS"
p.font.bold = True
p.font.size = Pt(12)
p.font.color.rgb = SIH_ORANGE

members = [
    ("1. Ayush Kumar", "Team Leader", "AI/ML, System Architecture, Overall Coordination", SIH_ORANGE),
    ("2. Prerana Mondal", "Team Member", "Frontend Development, UI/UX, Data Visualization", DARK_NAVY),
    ("3. Sondeep Kumar", "Team Member", "Backend Development, Database, API Integration", DARK_NAVY),
    ("4. Shinjini Lohar", "Team Member", "AI/ML, Computer Vision, Data Processing", DARK_NAVY),
    ("5. Subham Kumar Modi", "Team Member", "GIS Analysis, Mobile Application, QA Testing", DARK_NAVY),
    ("6. Rahul Deo", "Team Member", "Cloud Infrastructure, DevOps, Testing and Security", DARK_NAVY)
]

for name, role, resp, col in members:
    p = t_tf.add_paragraph()
    p.text = f"{name} "
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = col
    
    r = p.add_run()
    r.text = f"|  {role}"
    r.font.bold = True
    r.font.size = Pt(9.5)
    r.font.color.rgb = FOREST_GREEN if "Leader" in role else DARK_NAVY
    
    p_r = t_tf.add_paragraph()
    p_r.text = f"   Responsibilities: {resp}"
    p_r.font.size = Pt(9)
    p_r.font.color.rgb = TEXT_MUTED

# Slide 1 Footer
footer_left = slide1.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(9.5), Inches(0.32))
tf_fl = footer_left.text_frame
p_fl = tf_fl.paragraphs[0]
p_fl.text = "Team AlertNex | Smart India Hackathon 2026 | PS: SIH26001 | Disaster Management"
p_fl.font.size = Pt(9.5)
p_fl.font.color.rgb = TEXT_MUTED

footer_right = slide1.shapes.add_textbox(Inches(10.5), Inches(7.05), Inches(2.033), Inches(0.32))
tf_fr = footer_right.text_frame
p_fr = tf_fr.paragraphs[0]
p_fr.text = "Slide 1 of 6"
p_fr.font.bold = True
p_fr.font.size = Pt(9.5)
p_fr.font.color.rgb = DARK_NAVY
p_fr.alignment = PP_ALIGN.RIGHT

# ============================================================
# SLIDE 2: THE CHALLENGE
# ============================================================
slide2 = prs.slides.add_slide(blank_layout)
apply_clean_background(slide2)
add_slide_header_footer(slide2, 2, "THE CHALLENGE: Vulnerabilities in North Eastern Region")

# Problem Summary Banner
p_banner = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.6))
p_banner.fill.solid()
p_banner.fill.fore_color.rgb = RGBColor(0xFE, 0xF2, 0xF2)
p_banner.line.color.rgb = RISK_RED
p_banner.line.width = Pt(1.2)
pb_tf = p_banner.text_frame
pb_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
pb_tf.margin_left = Inches(0.2)
p = pb_tf.paragraphs[0]
p.text = "PROBLEM SUMMARY: "
p.font.bold = True
p.font.size = Pt(10.5)
p.font.color.rgb = RISK_RED
r = p.add_run()
r.text = "The North Eastern Region faces extreme precipitation and complex topography, causing sudden landslides that block lifeline highways, isolate tribal hamlets, and cut hospital access with delayed and purely reactive responses."
r.font.bold = False
r.font.color.rgb = TEXT_DARK

# 4 Key Problem Challenge Cards
challenges = [
    ("Heavy Rainfall & Slopes", "Cloudbursts (>150mm/day) oversaturate steep, fragile Himalayan terrain, causing sudden debris flows.", RISK_RED),
    ("Cut-off Lifeline Roads", "Single arterial highways (NH-10, NH-29) suffer recurring collapses, paralyzing state logistics.", RISK_ORANGE),
    ("Remote Village Isolation", "Deep-valley communities become completely isolated with no accessible route to emergency hospitals.", RISK_YELLOW),
    ("Zero Connectivity & Lag", "Severe storms knock out mobile towers; authorities receive ground reports hours or days too late.", DARK_NAVY)
]

c_w = Inches(2.78)
c_h = Inches(2.2)
for i, (title, desc, col) in enumerate(challenges):
    left = Inches(0.8 + i * 2.98)
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.15), c_w, c_h)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_LIGHT_GREY
    card.line.color.rgb = col
    card.line.width = Pt(1.5)
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.18)
    ctf.margin_top = Inches(0.15)
    
    p = ctf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = col
    
    p_desc = ctf.add_paragraph()
    p_desc.text = desc
    p_desc.font.size = Pt(9.5)
    p_desc.font.color.rgb = TEXT_DARK

# Visual Comparison: Current Approach vs Our Approach
comp_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.55), Inches(11.733), Inches(2.25))
comp_card.fill.solid()
comp_card.fill.fore_color.rgb = BG_LIGHT_GREY
comp_card.line.color.rgb = CARD_BORDER
comp_card.line.width = Pt(1.2)
comp_tf = comp_card.text_frame
comp_tf.word_wrap = True
comp_tf.margin_left = Inches(0.2)
comp_tf.margin_top = Inches(0.12)

p = comp_tf.paragraphs[0]
p.text = "VISUAL PARADIGM SHIFT: REACTIVE TO PROACTIVE"
p.font.bold = True
p.font.size = Pt(11.5)
p.font.color.rgb = DARK_NAVY

# Current Approach Box
curr_b = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(5.0), Inches(5.3), Inches(1.6))
curr_b.fill.solid()
curr_b.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF5)
curr_b.line.color.rgb = RISK_RED
curr_b.line.width = Pt(1.2)
cb_tf = curr_b.text_frame
cb_tf.word_wrap = True
cb_tf.margin_left = Inches(0.15)
cb_tf.margin_top = Inches(0.12)
p = cb_tf.paragraphs[0]
p.text = "CURRENT APPROACH (REACTIVE)"
p.font.bold = True
p.font.size = Pt(11)
p.font.color.rgb = RISK_RED

p = cb_tf.add_paragraph()
p.text = "Disaster Strikes  ➔  Detection  ➔  Delayed Response"
p.font.bold = True
p.font.size = Pt(10)
p.font.color.rgb = DARK_NAVY

p = cb_tf.add_paragraph()
p.text = "• Unaware of road blocks until stranded  |  No alternative route planning  |  Heavy casualty risk"
p.font.size = Pt(8.8)
p.font.color.rgb = TEXT_MUTED

# Our Approach Box
our_b = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(5.0), Inches(5.4), Inches(1.6))
our_b.fill.solid()
our_b.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
our_b.line.color.rgb = FOREST_GREEN
our_b.line.width = Pt(1.2)
ob_tf = our_b.text_frame
ob_tf.word_wrap = True
ob_tf.margin_left = Inches(0.15)
ob_tf.margin_top = Inches(0.12)
p = ob_tf.paragraphs[0]
p.text = "OUR APPROACH (TEAM ALERTNEX PROACTIVE)"
p.font.bold = True
p.font.size = Pt(11)
p.font.color.rgb = FOREST_GREEN

p = ob_tf.add_paragraph()
p.text = "Monitoring  ➔  Early Warning  ➔  Preparation  ➔  Faster Response"
p.font.bold = True
p.font.size = Pt(10)
p.font.color.rgb = DARK_NAVY

p = ob_tf.add_paragraph()
p.text = "• 6-12h pre-warning  |  Automated road blockage detection  |  Bypass routes active before collapse"
p.font.size = Pt(8.8)
p.font.color.rgb = TEXT_DARK

# ============================================================
# SLIDE 3: OUR PROPOSED SOLUTION
# ============================================================
slide3 = prs.slides.add_slide(blank_layout)
apply_clean_background(slide3)
add_slide_header_footer(slide3, 3, "OUR PROPOSED SOLUTION: Intelligent Disaster Monitoring Ecosystem")

# Left Column: Visual Ecosystem Diagram
eco_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(7.4), Inches(5.4))
eco_card.fill.solid()
eco_card.fill.fore_color.rgb = BG_LIGHT_GREY
eco_card.line.color.rgb = CARD_BORDER
eco_card.line.width = Pt(1.2)
e_tf = eco_card.text_frame
e_tf.word_wrap = True
e_tf.margin_left = Inches(0.2)
e_tf.margin_top = Inches(0.15)

p = e_tf.paragraphs[0]
p.text = "END-TO-END SOLUTION ECOSYSTEM FLOW"
p.font.bold = True
p.font.size = Pt(12)
p.font.color.rgb = DARK_NAVY

eco_steps = [
    ("1. MULTI-SOURCE DATA INGESTION", "Rainfall & Weather Data (IMD/GPM) + Soil Moisture (NASA SMAP) + Terrain/Slope (SRTM 30m DEM) + Satellite Imagery + Historical Landslide Inventory + Citizen & Field Reports", SIH_ORANGE),
    ("2. AI RISK ENGINE", "Machine learning ensemble analyzes multi-source environmental triggers, calculates slope stability index, and computes pixel-level susceptibility with Explainable AI factor weights.", DARK_NAVY),
    ("3. DYNAMIC LANDSLIDE RISK MAP", "Interactive GIS visualization (Leaflet/Mapbox + PostGIS) rendering 30m resolution dynamic risk zones across vulnerable North Eastern transportation corridors.", FOREST_GREEN),
    ("4. EARLY WARNING & ACTIONABLE RESPONSE", "Proactive alerts to SDRF/DDMA, village isolation predictions, automated emergency route suggestions, and offline mobile citizen guidance.", RISK_RED)
]

for title, desc, col in eco_steps:
    p = e_tf.add_paragraph()
    p.text = f"▼  {title}"
    p.font.bold = True
    p.font.size = Pt(10.2)
    p.font.color.rgb = col
    
    p_desc = e_tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.size = Pt(9)
    p_desc.font.color.rgb = TEXT_DARK

# Right Column: 4 Dynamic Risk Levels
risk_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.4), Inches(4.033), Inches(5.4))
risk_card.fill.solid()
risk_card.fill.fore_color.rgb = BG_LIGHT_GREY
risk_card.line.color.rgb = CARD_BORDER
risk_card.line.width = Pt(1.2)
r_tf = risk_card.text_frame
r_tf.word_wrap = True
r_tf.margin_left = Inches(0.2)
r_tf.margin_top = Inches(0.15)

p = r_tf.paragraphs[0]
p.text = "DYNAMIC RISK LEVELS"
p.font.bold = True
p.font.size = Pt(12)
p.font.color.rgb = DARK_NAVY

risks = [
    ("GREEN: LOW RISK (0 - 30%)", "Normal terrain stability. Routine satellite and weather polling. Regular traffic on mountain highways.", RISK_GREEN, RGBColor(0xEC, 0xFD, 0xF5)),
    ("YELLOW: MODERATE RISK (31 - 60%)", "Elevated moisture saturation or continuous rainfall. Advisory issued to highway patrols; field officers alerted to monitor slopes.", RISK_YELLOW, RGBColor(0xFF, 0xFB, 0xEB)),
    ("ORANGE: HIGH RISK (61 - 80%)", "Heavy precipitation + steep slope trigger. Pre-warning alerts to transport authorities; heavy freight restricted; alternative routes prepared.", RISK_ORANGE, RGBColor(0xFF, 0xF7, 0xED)),
    ("RED: CRITICAL RISK (81 - 100%)", "Imminent landslide danger. Instant Siren/SMS push to citizens, bypass emergency corridors activated, SDRF/NDRF teams pre-deployed.", RISK_RED, RGBColor(0xFE, 0xF2, 0xF2))
]

for title, desc, col, bg_col in risks:
    p = r_tf.add_paragraph()
    p.text = f"■  {title}"
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = col
    
    p_desc = r_tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.size = Pt(8.8)
    p_desc.font.color.rgb = TEXT_DARK

# ============================================================
# SLIDE 4: KEY INNOVATIONS
# ============================================================
slide4 = prs.slides.add_slide(blank_layout)
apply_clean_background(slide4)
add_slide_header_footer(slide4, 4, "KEY INNOVATIONS: Why Team AlertNex Solution is Different")

# 4 Innovation Cards in 2x2 Grid
innovations = [
    ("CARD 1: MULTI-SOURCE AI RISK ANALYSIS",
     "Combines Environmental & Ground-Level Information",
     [
         "Integrates static terrain data (slope, aspect, elevation) with dynamic live data (rainfall, soil saturation).",
         "Fuses satellite imagery with real-time crowdsourced citizen observations.",
         "Machine learning ensemble reduces false alarms and avoids alert fatigue."
     ], DARK_NAVY, Inches(0.8), Inches(1.4)),
     
    ("CARD 2: CONNECTIVITY IMPACT ANALYSIS",
     "The Core Innovation: Answering 'What Gets Cut Off?'",
     [
         "Road Blockage Analysis: Identifies exact highway segments threatened by impending slides.",
         "Village Isolation Analysis: Determines which remote hill villages lose vehicular access.",
         "Hospital Accessibility: Analyzes cut-off primary health centers and critical facilities.",
         "Alternative Routes: Autonomous graph routing calculates safe bypass emergency corridors."
     ], SIH_ORANGE, Inches(6.8), Inches(1.4)),
     
    ("CARD 3: OFFLINE COMMUNITY REPORTING",
     "Ground Intelligence in Zero-Internet Mountain Zones",
     [
         "Mobile application functions fully offline in remote zero-connectivity valleys.",
         "Citizens and field officers capture geotagged photos, cracks, and blocked roads.",
         "Reports are securely stored in local SQLite database on the mobile device.",
         "Automatically synchronizes with central server as soon as connection is restored."
     ], FOREST_GREEN, Inches(0.8), Inches(4.25)),
     
    ("CARD 4: EXPLAINABLE AI (XAI)",
     "Transparent Reasoning for Disaster Management Authorities",
     [
         "Explains WHY the risk is high rather than presenting an uninterpretable score.",
         "Concrete Factor Breakdown for Critical Alert (87%):",
         "  • Heavy Rainfall (24-hour accumulation): 42% contribution",
         "  • High Soil Moisture Saturation: 21% contribution",
         "  • Steep Terrain Gradient (>48°): 15% contribution",
         "  • Recent Citizen Crack Reports: 9% contribution"
     ], RISK_RED, Inches(6.8), Inches(4.25))
]

c_w = Inches(5.733)
c_h = Inches(2.65)
for title, subtitle, bullets, col, left, top in innovations:
    card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, c_w, c_h)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_LIGHT_GREY
    card.line.color.rgb = col
    card.line.width = Pt(1.5)
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.2)
    ctf.margin_top = Inches(0.12)
    
    p = ctf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = col
    
    p_sub = ctf.add_paragraph()
    p_sub.text = subtitle
    p_sub.font.italic = True
    p_sub.font.size = Pt(9.5)
    p_sub.font.color.rgb = DARK_NAVY
    
    for b in bullets:
        p_b = ctf.add_paragraph()
        p_b.text = f"• {b}"
        p_b.font.size = Pt(8.8)
        p_b.font.color.rgb = TEXT_DARK

# ============================================================
# SLIDE 5: SYSTEM ARCHITECTURE
# ============================================================
slide5 = prs.slides.add_slide(blank_layout)
apply_clean_background(slide5)
add_slide_header_footer(slide5, 5, "SYSTEM ARCHITECTURE: 4-Tier Flow & Technology Stack")

# 4 Architecture Flow Layers
arch_layers = [
    ("LAYER 1: DATA COLLECTION", 
     "Weather & Rainfall (IMD API, GPM) • Soil Moisture (NASA SMAP) • Terrain & Elevation (SRTM 30m DEM) • Satellite Data (Sentinel-2) • Citizen & Field Officer Reports",
     SIH_ORANGE),
    ("LAYER 2: AI AND ANALYTICS ENGINE",
     "Data Processing (Pandas, NumPy) • Machine Learning Risk Analysis (Scikit-learn, XGBoost) • Explainable AI Module (SHAP Factor Breakdown)",
     DARK_NAVY),
    ("LAYER 3: GIS AND IMPACT ANALYSIS",
     "PostgreSQL + PostGIS Spatial Geodatabase • Dynamic Risk Mapping • Road Blockage Analysis • Village Isolation Detection • Emergency Alternative Routes (Dijkstra)",
     FOREST_GREEN),
    ("LAYER 4: OUTPUT & ACTION",
     "Authority Command Dashboard (React.js + Leaflet/Mapbox) • Mobile Application (Flutter + SQLite) • Early Warning Alerts (Firebase Push & SMS)",
     RISK_RED)
]

for i, (l_title, l_desc, l_col) in enumerate(arch_layers):
    top = Inches(1.4 + i * 1.15)
    layer_b = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.0))
    layer_b.fill.solid()
    layer_b.fill.fore_color.rgb = BG_LIGHT_GREY
    layer_b.line.color.rgb = l_col
    layer_b.line.width = Pt(1.5)
    ltf = layer_b.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0.2)
    ltf.margin_top = Inches(0.1)
    
    p = ltf.paragraphs[0]
    p.text = l_title
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = l_col
    
    p_desc = ltf.add_paragraph()
    p_desc.text = l_desc
    p_desc.font.size = Pt(9.2)
    p_desc.font.color.rgb = TEXT_DARK

# Technology Stack Box at Bottom
tech_b = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.05), Inches(11.733), Inches(0.88))
tech_b.fill.solid()
tech_b.fill.fore_color.rgb = RGBColor(0x0A, 0x19, 0x2F)
tech_b.line.color.rgb = DARK_NAVY
tech_b.line.width = Pt(1)
ttf = tech_b.text_frame
ttf.word_wrap = True
ttf.margin_left = Inches(0.2)
ttf.margin_top = Inches(0.08)

p = ttf.paragraphs[0]
p.text = "PURPOSE-BUILT TECHNOLOGY STACK:"
p.font.bold = True
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(0x64, 0xFF, 0xDA)

p_stack = ttf.add_paragraph()
p_stack.text = "React.js (Dashboard)  •  Flutter (Mobile)  •  Python + FastAPI (Backend)  •  Scikit-learn (AI/ML)  •  PostgreSQL + PostGIS (Spatial DB)  •  Leaflet (GIS Mapping)  •  SQLite (Offline Mobile DB)  •  Firebase (Cloud Alerts)"
p_stack.font.size = Pt(9.2)
p_stack.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ============================================================
# SLIDE 6: IMPACT, FEASIBILITY & FUTURE
# ============================================================
slide6 = prs.slides.add_slide(blank_layout)
apply_clean_background(slide6)
add_slide_header_footer(slide6, 6, "IMPACT, FEASIBILITY & FUTURE: The AlertNex Roadmap")

# 3 Columns
columns = [
    ("EXPECTED IMPACT", [
        "Earlier disaster preparedness with 6-12 hour warning windows.",
        "Faster emergency response and rescue mobilization for SDRF & BRO.",
        "Better protection for vulnerable remote tribal hill communities.",
        "Reduced village isolation through proactive road clearing.",
        "Improved data-driven decision-making for district magistrates.",
        "Better infrastructure planning and reduced recurring economic losses."
    ], FOREST_GREEN, Inches(0.8)),
    
    ("FEASIBILITY", [
        "Uses open environmental and geographical data (IMD, NASA, SRTM).",
        "Can start with selected high-risk pilot areas (e.g. Sikkim or Assam).",
        "Modular system architecture ensures seamless independent progress.",
        "Realistic student-level prototype deployable on cloud free-tiers.",
        "Uses reliable open-source technologies (FastAPI, PostGIS, Flutter).",
        "Operates without requiring crores in expensive physical hardware sensors."
    ], DARK_NAVY, Inches(4.75)),
    
    ("FUTURE SCOPE", [
        "Drone monitoring for automated aerial slope crack verification.",
        "Advanced satellite radar (InSAR) monitoring for millimeter ground movement.",
        "Regional language support (Assamese, Bengali, Nepali, Mizo, Khasi).",
        "Emergency service integration directly with 112 India and BRO Swastik.",
        "Expansion to other Himalayan regions (Uttarakhand, Himachal Pradesh)."
    ], SIH_ORANGE, Inches(8.7))
]

c_w = Inches(3.8)
c_h = Inches(4.6)
for title, points, col, left in columns:
    card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), c_w, c_h)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_LIGHT_GREY
    card.line.color.rgb = col
    card.line.width = Pt(1.5)
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.18)
    ctf.margin_top = Inches(0.15)
    
    p = ctf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(11.5)
    p.font.color.rgb = col
    
    for pt in points:
        p_pt = ctf.add_paragraph()
        p_pt.text = f"• {pt}"
        p_pt.font.size = Pt(8.8)
        p_pt.font.color.rgb = TEXT_DARK

# Final Message Banner at Bottom
fm_banner = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.15), Inches(11.733), Inches(0.75))
fm_banner.fill.solid()
fm_banner.fill.fore_color.rgb = RGBColor(0x0A, 0x19, 0x2F)
fm_banner.line.color.rgb = SIH_ORANGE
fm_banner.line.width = Pt(1.5)
fm_tf = fm_banner.text_frame
fm_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = fm_tf.paragraphs[0]
p.text = "TEAM ALERTNEX  •  AI-POWERED EARLY WARNING FOR SAFER COMMUNITIES"
p.font.bold = True
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = "SIH2026_AlertNex_Official_6Slides.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to: {output_path}")

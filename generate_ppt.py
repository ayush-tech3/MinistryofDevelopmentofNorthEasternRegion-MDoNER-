"""
SIH 2026 PPT Generator for Team AlertNex
Problem Statement: AI-Based Early Warning and Landslide Risk Monitoring System in NER
Problem Statement ID: SH2001 | Organization: Ministry of Development of North Eastern Region (MDoNER)
Theme: Disaster Management | Category: Software Edition

Generates:
1. SIH2026_AlertNex_Presentation.pptx (Strict Official SIH 6-Slide Submission Presentation)
2. SIH2026_AlertNex_Extended_10Slides.pptx (Extended 10-Slide Presentation for Internal Rounds)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
DARK_NAVY = RGBColor(0x0A, 0x1A, 0x2F)
DEEP_BLUE = RGBColor(0x0D, 0x27, 0x47)
ACCENT_BLUE = RGBColor(0x00, 0x7B, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
ACCENT_GREEN = RGBColor(0x00, 0xC9, 0x7B)
ACCENT_RED = RGBColor(0xFF, 0x4D, 0x4D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GOLD = RGBColor(0xFF, 0xC1, 0x07)
LIGHT_BLUE_BG = RGBColor(0x0F, 0x2B, 0x4C)
CARD_BORDER = RGBColor(0x1A, 0x3A, 0x5C)

# ── Slide dimensions (Widescreen 16:9) ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Official 6-Member Team Mandate (SIH Rule Compliant) ──
TEAM_MEMBERS = [
    ("👤 Ayush Kumar", "Team Leader", "Full-Stack Dev, AI/ML, System Architecture", GOLD),
    ("👤 Prerana Mondal", "Member (Female)", "Frontend Dev, UI/UX Design, Data Visualization", ACCENT_BLUE),
    ("👤 Sondeep Kumar", "Member", "Backend Dev, PostgreSQL/PostGIS, REST APIs", ACCENT_CYAN),
    ("👤 Harshit Jha", "Member", "AI/ML Engineer, Computer Vision, Predictive Models", ACCENT_GREEN),
    ("👤 Saksham Singh", "Member", "Geospatial & GIS Analyst, Mobile PWA, QA & Testing", ACCENT_ORANGE),
    ("👤 [Member 6 Name]", "Member (To Be Updated)", "Cloud Infrastructure, Edge Computing & DevOps Security", ACCENT_RED),
]

METRICS = [
    ("60%+", "Casualty Reduction", "2 to 6-hour proactive warning window enables safe mass evacuation before collapse.", ACCENT_GREEN),
    ("₹500 Cr+", "Annual Savings", "Mitigates highway transit disruption, equipment loss & post-disaster rescue bills.", ACCENT_BLUE),
    ("45M+", "Citizens Protected", "Comprehensive hazard shielding across all 8 North Eastern states and border villages.", ACCENT_ORANGE),
    (">90%", "Prediction Precision", "Multi-source AI ensemble eliminates false alarms and mitigates alert fatigue.", ACCENT_CYAN),
]

SDGS = [
    ("SDG 11: Sustainable Cities", "Target 11.5: Protect vulnerable hill settlements & reduce disaster casualties.", ACCENT_ORANGE),
    ("SDG 13: Climate Action", "Target 13.1: Adaptive resilience against climate-driven extreme cloudbursts.", ACCENT_GREEN),
    ("SDG 9: Resilient Infrastructure", "Target 9.1: Proactively safeguard highway lifelines (NH-10, NH-29) and bridges.", ACCENT_BLUE),
    ("SDG 17: Multi-Agency Partnerships", "Target 17.16: Unified coordination between MDoNER, NDMA, IMD, GSI & BRO.", ACCENT_CYAN),
]


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a shape to the slide."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14, font_color=WHITE,
                bold=False, italic=False, alignment=PP_ALIGN.LEFT, font_name='Calibri',
                vertical_anchor=MSO_ANCHOR.TOP):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, texts, font_name='Calibri',
                   vertical_anchor=MSO_ANCHOR.TOP):
    """Add a text box with multiple formatted paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, t in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = t.get('text', '')
        p.font.size = Pt(t.get('font_size', 14))
        p.font.color.rgb = t.get('font_color', WHITE)
        p.font.bold = t.get('bold', False)
        p.font.italic = t.get('italic', False)
        p.font.name = font_name
        p.alignment = t.get('alignment', PP_ALIGN.LEFT)
        if 'space_after' in t:
            p.space_after = Pt(t['space_after'])
        if 'space_before' in t:
            p.space_before = Pt(t['space_before'])
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_CYAN, height=Pt(3)):
    """Add an accent line."""
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill_color=color)


def add_card(slide, left, top, width, height, fill_color=LIGHT_BLUE_BG, line_color=CARD_BORDER, line_width=1):
    """Add a styled rounded rectangle card."""
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                     fill_color=fill_color, line_color=line_color, line_width=line_width)


def add_footer_bar(slide, slide_num, total=6):
    """Add standardized official footer bar."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05),
              SLIDE_WIDTH, Inches(0.45), fill_color=RGBColor(0x05, 0x10, 0x20))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.03),
              SLIDE_WIDTH, Pt(2), fill_color=ACCENT_CYAN)
    add_textbox(slide, Inches(0.5), Inches(7.08), Inches(8), Inches(0.35),
                "Team AlertNex  |  Smart India Hackathon 2026  |  PS: SH2001 (MDoNER)  |  Theme: Disaster Management",
                font_size=9, font_color=MEDIUM_GRAY)
    add_textbox(slide, Inches(10.5), Inches(7.08), Inches(2.3), Inches(0.35),
                f"Slide {slide_num} of {total}", font_size=9, font_color=MEDIUM_GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_slide_header(slide, title, category="SMART INDIA HACKATHON 2026  |  MDoNER", accent_color=ACCENT_CYAN):
    """Add consistent slide header."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
              SLIDE_WIDTH, Inches(0.06), fill_color=accent_color)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.3),
                category, font_size=10, font_color=ACCENT_ORANGE, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.45), Inches(10), Inches(0.55),
                title, font_size=24, font_color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.5), Inches(1.05), Inches(3.5), color=accent_color)


# ==============================================================================
# 1. OFFICIAL 6-SLIDE PRESENTATION BUILDER (STRICT SIH SUBMISSION FORMAT)
# ==============================================================================
def create_official_6slide_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 1: TITLE SLIDE (Official SIH Template Standard)
    # ──────────────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1, DARK_NAVY)

    # Top accent bar
    add_shape(s1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), fill_color=ACCENT_CYAN)
    add_shape(s1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.06), SLIDE_HEIGHT, fill_color=ACCENT_BLUE)

    # Badges row
    b1 = add_shape(s1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.35), Inches(3.2), Inches(0.45),
                   fill_color=RGBColor(0x0D, 0x47, 0xA1), line_color=ACCENT_BLUE, line_width=1)
    b1.text_frame.paragraphs[0].text = "SMART INDIA HACKATHON 2026"
    b1.text_frame.paragraphs[0].font.size = Pt(12)
    b1.text_frame.paragraphs[0].font.color.rgb = WHITE
    b1.text_frame.paragraphs[0].font.bold = True
    b1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    b2 = add_shape(s1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.85), Inches(0.35), Inches(1.8), Inches(0.45),
                   fill_color=ACCENT_ORANGE)
    b2.text_frame.paragraphs[0].text = "PS ID: SH2001"
    b2.text_frame.paragraphs[0].font.size = Pt(12)
    b2.text_frame.paragraphs[0].font.color.rgb = WHITE
    b2.text_frame.paragraphs[0].font.bold = True
    b2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    b3 = add_shape(s1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(0.35), Inches(2.0), Inches(0.45),
                   fill_color=ACCENT_GREEN)
    b3.text_frame.paragraphs[0].text = "Software Edition"
    b3.text_frame.paragraphs[0].font.size = Pt(12)
    b3.text_frame.paragraphs[0].font.color.rgb = WHITE
    b3.text_frame.paragraphs[0].font.bold = True
    b3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    b4 = add_shape(s1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.95), Inches(0.35), Inches(2.4), Inches(0.45),
                   fill_color=RGBColor(0x13, 0x88, 0x08))
    b4.text_frame.paragraphs[0].text = "Theme: Disaster Management"
    b4.text_frame.paragraphs[0].font.size = Pt(11)
    b4.text_frame.paragraphs[0].font.color.rgb = WHITE
    b4.text_frame.paragraphs[0].font.bold = True
    b4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Main Title
    add_multi_text(s1, Inches(0.5), Inches(1.0), Inches(8.5), Inches(2.2), [
        {'text': 'AI-Based Early Warning &', 'font_size': 32, 'font_color': WHITE, 'bold': True, 'space_after': 2},
        {'text': 'Landslide Risk Monitoring System', 'font_size': 32, 'font_color': ACCENT_CYAN, 'bold': True, 'space_after': 4},
        {'text': 'for North Eastern Region (NER) of India', 'font_size': 20, 'font_color': LIGHT_GRAY, 'space_after': 6},
    ])

    # Right side PS & Institute Details Card
    info_card = add_card(s1, Inches(9.2), Inches(0.95), Inches(3.6), Inches(2.3))
    add_accent_line(s1, Inches(9.3), Inches(1.05), Inches(3.4), color=ACCENT_CYAN)
    add_textbox(s1, Inches(9.4), Inches(1.15), Inches(3.2), Inches(0.3),
                "OFFICIAL PROBLEM DETAILS", font_size=11, font_color=ACCENT_CYAN, bold=True)
    ps_info = [
        ("Organization:", "Ministry of Development of NER (MDoNER)"),
        ("Problem Statement:", "SH2001 (Software Edition)"),
        ("Target Region:", "8 North Eastern States"),
        ("Institute Name:", "[Your College / Institute Name, City]"),
        ("Faculty Mentor:", "[Faculty Mentor Name / Designation]"),
    ]
    for i, (k, v) in enumerate(ps_info):
        add_textbox(s1, Inches(9.4), Inches(1.5 + i * 0.32), Inches(1.3), Inches(0.3),
                    k, font_size=9, font_color=MEDIUM_GRAY)
        add_textbox(s1, Inches(10.6), Inches(1.5 + i * 0.32), Inches(2.1), Inches(0.3),
                    v, font_size=9, font_color=WHITE, bold=True)

    # Team Header & Tagline
    add_accent_line(s1, Inches(0.5), Inches(3.25), Inches(4.5), color=ACCENT_ORANGE)
    add_textbox(s1, Inches(0.5), Inches(3.35), Inches(5.5), Inches(0.4),
                "Team AlertNex  (Lead: Ayush Kumar)", font_size=20, font_color=ACCENT_ORANGE, bold=True)
    add_textbox(s1, Inches(0.5), Inches(3.75), Inches(12.3), Inches(0.35),
                "\"Turning raw environmental data into actionable intelligence for landslide disaster resilience in NER\"",
                font_size=12, font_color=ACCENT_GREEN, italic=True)

    # Team Members Card (6 Members Mandate Compliant)
    team_card = add_card(s1, Inches(0.5), Inches(4.2), Inches(12.33), Inches(2.65))
    add_accent_line(s1, Inches(0.6), Inches(4.28), Inches(12.1), color=ACCENT_CYAN)
    add_textbox(s1, Inches(0.7), Inches(4.35), Inches(6), Inches(0.3),
                "TEAM MEMBERS  (SIH 6-Member Mandate Compliant)", font_size=11, font_color=ACCENT_CYAN, bold=True)
    add_textbox(s1, Inches(7.0), Inches(4.35), Inches(5.6), Inches(0.3),
                "✔ Female Representation: Prerana Mondal  |  100% SIH Rule Compliant",
                font_size=10, font_color=GOLD, italic=True, alignment=PP_ALIGN.RIGHT)

    for i, (name, role, skills, color) in enumerate(TEAM_MEMBERS):
        x = Inches(0.65) + Inches(i * 2.02)
        # Inner card
        add_shape(s1, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.75), Inches(1.95), Inches(1.95),
                  fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
        add_shape(s1, MSO_SHAPE.RECTANGLE, x + Inches(0.05), Inches(4.78), Inches(1.85), Pt(2), fill_color=color)

        # Avatar circle
        av = add_shape(s1, MSO_SHAPE.OVAL, x + Inches(0.65), Inches(4.88), Inches(0.65), Inches(0.65), fill_color=color)
        av.text_frame.paragraphs[0].text = name[2] if name[2] != '[' else '6'
        av.text_frame.paragraphs[0].font.size = Pt(16)
        av.text_frame.paragraphs[0].font.color.rgb = WHITE
        av.text_frame.paragraphs[0].font.bold = True
        av.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(s1, x + Inches(0.05), Inches(5.6), Inches(1.85), Inches(0.3),
                    name.replace("👤 ", ""), font_size=10, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s1, x + Inches(0.05), Inches(5.88), Inches(1.85), Inches(0.25),
                    role, font_size=8.5, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s1, x + Inches(0.05), Inches(6.12), Inches(1.85), Inches(0.55),
                    skills, font_size=7.5, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_footer_bar(s1, 1, 6)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 2: IDEA TITLE & PROPOSED SOLUTION (Problem Context & Solution)
    # ──────────────────────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2, DARK_NAVY)
    add_slide_header(s2, "IDEA TITLE & PROPOSED SOLUTION", "SMART INDIA HACKATHON 2026  |  PROBLEM UNDERSTANDING & SOLUTION", ACCENT_GREEN)

    # Title Card
    idea_box = add_card(s2, Inches(0.5), Inches(1.15), Inches(12.33), Inches(0.55), fill_color=RGBColor(0x00, 0x2B, 0x1A), line_color=ACCENT_GREEN)
    add_textbox(s2, Inches(0.6), Inches(1.22), Inches(12.1), Inches(0.4),
                "AlertNex: AI-Powered Landslide Early Warning & Risk Monitoring Platform for NER",
                font_size=15, font_color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # Left Column: Problem Background & Context (Why MDoNER needs this)
    add_card(s2, Inches(0.5), Inches(1.85), Inches(4.5), Inches(4.2))
    add_accent_line(s2, Inches(0.6), Inches(1.95), Inches(4.3), color=ACCENT_ORANGE)
    add_textbox(s2, Inches(0.7), Inches(2.05), Inches(4.1), Inches(0.3),
                "BACKGROUND & REGIONAL CONTEXT", font_size=12, font_color=ACCENT_ORANGE, bold=True)

    bg_points = [
        ("Fragile Geology:", "Young, seismically active Himalayan & Indo-Burman mountain ranges spanning all 8 NER states."),
        ("Extreme Monsoons:", "Torrential rainfall (>2,000 mm/year) triggers sudden debris flows, slope failures, and flash floods."),
        ("Severe Toll:", "4,796+ recorded landslides (2015-2024), ~1,200 annual casualties, and ₹500+ Cr annual economic losses."),
        ("Transport Isolation:", "Vital highway lifelines NH-10 (Sikkim) and NH-29 (Nagaland/Manipur) repeatedly severed for weeks."),
        ("Status Quo Gaps:", "Current monitoring is 70%+ reactive post-disaster, manual, single-sensor prone, and fails in hill telecom blackouts."),
    ]
    for i, (title, desc) in enumerate(bg_points):
        y = Inches(2.4 + i * 0.65)
        add_textbox(s2, Inches(0.7), y, Inches(4.1), Inches(0.22),
                    f"▸ {title}", font_size=10, font_color=WHITE, bold=True)
        add_textbox(s2, Inches(0.85), y + Inches(0.22), Inches(3.95), Inches(0.4),
                    desc, font_size=8.5, font_color=LIGHT_GRAY)

    # Right Column: Proposed Solution Pillars (6 Core Pillars)
    add_card(s2, Inches(5.15), Inches(1.85), Inches(7.68), Inches(4.2))
    add_accent_line(s2, Inches(5.25), Inches(1.95), Inches(7.48), color=ACCENT_CYAN)
    add_textbox(s2, Inches(5.35), Inches(2.05), Inches(7.2), Inches(0.3),
                "ALERTNEX ARCHITECTURAL SOLUTION PILLARS", font_size=12, font_color=ACCENT_CYAN, bold=True)

    pillars = [
        ("Multi-Source Ingestion", "Automates live ingestion of IMD weather radar, Sentinel-2 optical/SAR satellite imagery, GSI lithology, and ground IoT soil sensors into unified spatial pipeline.", ACCENT_BLUE),
        ("Hybrid AI/ML Ensemble", "Fuses Random Forest, XGBoost, CNN displacement detection, and LSTM rainfall forecasting to predict slope collapse 2-6 hours prior (>90% accuracy).", ACCENT_CYAN),
        ("GIS Risk Heatmap & Evacuation", "Interactive Leaflet & PostGIS visualizer mapping live hazard danger zones, road blockages, and automated safe evacuation corridors.", ACCENT_GREEN),
        ("Multi-Channel Local Alerts", "Instant warning broadcast via SMS, siren triggers, and synthesized voice calls in native NER languages (Assamese, Mizo, Manipuri, Khasi, Nagamese).", ACCENT_ORANGE),
        ("Field Crowdsourcing (PWA)", "Progressive Web App empowering citizens and field workers to upload geo-tagged images of ground fissures and water seepage to calibrate AI models.", GOLD),
        ("Offline Edge & Low-Network Mode", "Local edge caching with store-and-forward sync plus automatic SMS fallback ensures 100% operation during hill network blackouts.", ACCENT_RED),
    ]

    for i, (p_title, p_desc, p_col) in enumerate(pillars):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(5.35) + Inches(col_idx * 3.7)
        y = Inches(2.4) + Inches(row_idx * 1.15)

        add_shape(s2, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.55), Inches(1.05),
                  fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
        add_shape(s2, MSO_SHAPE.RECTANGLE, x + Inches(0.05), y + Inches(0.05), Inches(3.45), Pt(2), fill_color=p_col)
        add_textbox(s2, x + Inches(0.1), y + Inches(0.1), Inches(3.35), Inches(0.25),
                    p_title, font_size=10, font_color=p_col, bold=True)
        add_textbox(s2, x + Inches(0.1), y + Inches(0.35), Inches(3.35), Inches(0.65),
                    p_desc, font_size=8, font_color=LIGHT_GRAY)

    # Bottom Novelty Callout Bar
    add_shape(s2, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.15), Inches(12.33), Inches(0.75),
              fill_color=RGBColor(0x1A, 0x2D, 0x45), line_color=ACCENT_CYAN, line_width=1)
    add_textbox(s2, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.25),
                "⭐ CORE NOVELTY & COMPETITIVE ADVANTAGE", font_size=10, font_color=GOLD, bold=True)
    add_textbox(s2, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.4),
                "1) Multi-modal AI fusing spaceborne InSAR, weather radar & ground sensors.  2) 100% Offline-First design overcoming hill telecom loss.  "
                "3) Vernacular voice synthesis bridging ethnic linguistic barriers.  4) Closed-loop crowdsourcing validating regional ground-truth.",
                font_size=8.5, font_color=WHITE)

    add_footer_bar(s2, 2, 6)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 3: TECHNICAL APPROACH & ARCHITECTURE
    # ──────────────────────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3, DARK_NAVY)
    add_slide_header(s3, "TECHNICAL APPROACH & SYSTEM ARCHITECTURE", "SMART INDIA HACKATHON 2026  |  TECHNICAL IMPLEMENTATION", ACCENT_BLUE)

    # 4 Architecture Layers
    layers = [
        ("LAYER 1: DATA INGESTION & SENSING", ACCENT_BLUE, [
            "Satellite Feeds: Sentinel-2 Multispectral & Landsat-8/9 Surface Temperature",
            "Mausam Weather Radar: Live IMD Rainfall & AWS precipitation APIs",
            "Geological & Elevation: GSI Lithology fault lines & SRTM 30m DEM slope models",
            "NRSC Historical Catalog: ISRO Bhuvan historical landslide spatial points",
            "Community Inputs: Geo-tagged crack photos & IoT soil moisture/pore sensors",
        ]),
        ("LAYER 2: AI/ML PREDICTIVE ENGINE", ACCENT_CYAN, [
            "Spatial Feature Engineering: NDVI vegetation index, slope aspect & curvature",
            "Susceptibility Classification: Random Forest & XGBoost ensemble models",
            "Surface Deformation Detection: CNN analyzing optical displacement patterns",
            "Precipitation Threshold Trigger: LSTM neural network on rainfall time-series",
            "Dynamic Probability Scoring: >90% precision with automated drift re-calibration",
        ]),
        ("LAYER 3: GEOSPATIAL & ALERT SERVICES", ACCENT_GREEN, [
            "PostGIS Database: Spatio-temporal relational storage for rapid polygon queries",
            "GeoServer & Vector Tiles: Real-time dynamic risk heatmaps generation",
            "Emergency Routing Engine: OpenStreetMap shortest safe corridor pathfinding",
            "Multi-Channel Dispatch: Twilio SMS, Firebase FCM push & automated siren relays",
            "Multilingual NLP: Text & speech synthesis in 10+ North Eastern languages",
        ]),
        ("LAYER 4: USER APPLICATIONS & APIS", ACCENT_ORANGE, [
            "Command Console: React.js web dashboard for District Admins & SDMA officers",
            "Field Personnel PWA: Offline-first mobile app with local SQLite/IndexedDB cache",
            "Public Travel Portal: Real-time highway status & landslide travel advisories",
            "External REST APIs: Interoperability with NDMA national portal & BRO command",
        ]),
    ]

    for i, (l_title, l_col, l_items) in enumerate(layers):
        x = Inches(0.5) + Inches(i * 3.1)
        add_card(s3, x, Inches(1.2), Inches(2.95), Inches(3.6))
        add_shape(s3, MSO_SHAPE.RECTANGLE, x + Inches(0.05), Inches(1.25), Inches(2.85), Pt(3), fill_color=l_col)

        add_textbox(s3, x + Inches(0.1), Inches(1.35), Inches(2.75), Inches(0.4),
                    l_title, font_size=10, font_color=l_col, bold=True, alignment=PP_ALIGN.CENTER)

        for j, item in enumerate(l_items):
            add_textbox(s3, x + Inches(0.12), Inches(1.8 + j * 0.58), Inches(2.7), Inches(0.55),
                        f"• {item}", font_size=7.5, font_color=LIGHT_GRAY)

    # Technology Stack & Pipeline Summary (Bottom Half)
    add_card(s3, Inches(0.5), Inches(4.9), Inches(12.33), Inches(2.0))
    add_accent_line(s3, Inches(0.6), Inches(4.98), Inches(12.1), color=ACCENT_CYAN)
    add_textbox(s3, Inches(0.7), Inches(5.05), Inches(6), Inches(0.25),
                "COMPREHENSIVE TECHNOLOGY STACK & DATA PIPELINE", font_size=11, font_color=ACCENT_CYAN, bold=True)

    techs = [
        ("AI / ML & Data Science", "Python, TensorFlow, PyTorch, scikit-learn, XGBoost, OpenCV, Pandas, NumPy", ACCENT_BLUE),
        ("Backend & Databases", "FastAPI (Python), Node.js / Express, PostgreSQL + PostGIS, Redis Cache", ACCENT_CYAN),
        ("Frontend & Geospatial", "React.js, MapLibre GL, Leaflet.js, GeoServer, Tailwind CSS, PWA Service Workers", ACCENT_GREEN),
        ("Cloud, Edge & Alerts", "Docker, Kubernetes, AWS/GCP, Nginx, Twilio SMS API, Firebase FCM, Webhooks", ACCENT_ORANGE),
    ]
    for i, (t_cat, t_tools, t_col) in enumerate(techs):
        x = Inches(0.7) + Inches(i * 3.05)
        add_shape(s3, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.35), Inches(2.9), Inches(0.75),
                  fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
        add_textbox(s3, x + Inches(0.1), Inches(5.4), Inches(2.7), Inches(0.25),
                    t_cat, font_size=9, font_color=t_col, bold=True)
        add_textbox(s3, x + Inches(0.1), Inches(5.65), Inches(2.7), Inches(0.4),
                    t_tools, font_size=7.5, font_color=LIGHT_GRAY)

    # Pipeline Flow Steps
    add_textbox(s3, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.6),
                "Pipeline Execution: 1. Ingest (IMD/Sat Feeds) ➔ 2. Feature Extraction (Slope/Rainfall) ➔ "
                "3. Ensemble Inference (Risk > 90%) ➔ 4. Alert Broadcast (<30s) ➔ 5. Evacuation Routing",
                font_size=9, font_color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    add_footer_bar(s3, 3, 6)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 4: FEASIBILITY & VIABILITY
    # ──────────────────────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4, DARK_NAVY)
    add_slide_header(s4, "FEASIBILITY & VIABILITY", "SMART INDIA HACKATHON 2026  |  IMPLEMENTATION & RISK MITIGATION", ACCENT_GREEN)

    # Left Column: 36-Hour Hackathon Execution Plan
    add_card(s4, Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.7))
    add_accent_line(s4, Inches(0.6), Inches(1.3), Inches(4.3), color=ACCENT_GREEN)
    add_textbox(s4, Inches(0.7), Inches(1.4), Inches(4.1), Inches(0.3),
                "36-HOUR HACKATHON SPRINT PLAN", font_size=12, font_color=ACCENT_GREEN, bold=True)

    sprints = [
        ("Hours 00 - 08", "Environment & ML Baseline", "Docker stack setup, NRSC historical data ingestion, baseline Random Forest & XGBoost model training (>85% initial accuracy).", ACCENT_BLUE),
        ("Hours 08 - 18", "GIS Backend & Spatial Engine", "PostGIS spatial schema, GeoServer vector tile rendering, dynamic hazard contour generation, and REST APIs.", ACCENT_CYAN),
        ("Hours 18 - 28", "UI Dashboard & Alert Gateway", "React command console with interactive Leaflet map, PWA offline caching service workers, Twilio SMS broadcast pipeline.", ACCENT_GREEN),
        ("Hours 28 - 36", "Scenario Testing & Polishing", "End-to-end simulated Sikkim/Assam monsoon flood event, edge network failure failover test, bug fixing & demo dry runs.", ACCENT_ORANGE),
    ]
    for i, (hrs, title, desc, col) in enumerate(sprints):
        y = Inches(1.75 + i * 1.2)
        add_shape(s4, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(4.1), Inches(1.1),
                  fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
        add_shape(s4, MSO_SHAPE.RECTANGLE, Inches(0.75), y + Inches(0.05), Inches(4.0), Pt(2), fill_color=col)
        add_textbox(s4, Inches(0.85), y + Inches(0.1), Inches(1.4), Inches(0.25),
                    hrs, font_size=9, font_color=col, bold=True)
        add_textbox(s4, Inches(2.25), y + Inches(0.1), Inches(2.4), Inches(0.25),
                    title, font_size=9.5, font_color=WHITE, bold=True)
        add_textbox(s4, Inches(0.85), y + Inches(0.38), Inches(3.8), Inches(0.65),
                    desc, font_size=8, font_color=LIGHT_GRAY)

    # Right Column Top: Scalability Roadmap
    add_card(s4, Inches(5.15), Inches(1.2), Inches(7.68), Inches(2.55))
    add_accent_line(s4, Inches(5.25), Inches(1.3), Inches(7.48), color=ACCENT_ORANGE)
    add_textbox(s4, Inches(5.35), Inches(1.4), Inches(7.2), Inches(0.3),
                "POST-HACKATHON SCALABILITY ROADMAP", font_size=12, font_color=ACCENT_ORANGE, bold=True)

    roadmap = [
        ("Phase 1: Hackathon", "36 Hours", "Working prototype on 2 pilot NER districts with live IMD weather feed.", ACCENT_BLUE),
        ("Phase 2: On-Ground Pilot", "3 Months", "Deployment across 5 high-risk districts in Sikkim & Assam with ground IoT nodes.", ACCENT_CYAN),
        ("Phase 3: NER-Wide Rollout", "6 Months", "Full coverage across all 8 NER states integrating State Disaster Management Authorities.", ACCENT_GREEN),
        ("Phase 4: National Integration", "12 Months", "Pan-India integration with NDMA national portal, BRO, and Indian Railways.", ACCENT_ORANGE),
    ]
    for i, (p_name, dur, p_desc, p_col) in enumerate(roadmap):
        x = Inches(5.35) + Inches(i * 1.85)
        add_shape(s4, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.75), Inches(1.75), Inches(1.85),
                  fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
        add_shape(s4, MSO_SHAPE.RECTANGLE, x + Inches(0.05), Inches(1.8), Inches(1.65), Pt(2), fill_color=p_col)
        add_textbox(s4, x + Inches(0.08), Inches(1.88), Inches(1.6), Inches(0.3),
                    p_name, font_size=8.5, font_color=p_col, bold=True)
        add_textbox(s4, x + Inches(0.08), Inches(2.2), Inches(1.6), Inches(0.25),
                    dur, font_size=8, font_color=GOLD, bold=True)
        add_textbox(s4, x + Inches(0.08), Inches(2.45), Inches(1.6), Inches(1.1),
                    p_desc, font_size=7.5, font_color=LIGHT_GRAY)

    # Right Column Bottom: Risk Analysis & Mitigations
    add_card(s4, Inches(5.15), Inches(3.85), Inches(7.68), Inches(3.05))
    add_accent_line(s4, Inches(5.25), Inches(3.95), Inches(7.48), color=ACCENT_RED)
    add_textbox(s4, Inches(5.35), Inches(4.05), Inches(7.2), Inches(0.3),
                "ENGINEERED RISK ANALYSIS & MITIGATION MATRIX", font_size=12, font_color=ACCENT_RED, bold=True)

    risks = [
        ("Hill Telecom Blackouts", "Delays early warning", "Offline-first PWA caching + local mesh + automated SMS broadcast fallback", "RESOLVED"),
        ("Complex Topography", "False alarms / errors", "Ensemble AI fusing local rain gauges, InSAR satellite data & terrain DEM", "RESOLVED"),
        ("15+ Native Dialects", "Comprehension barrier", "Automated multilingual text & speech synthesis in Assamese, Mizo, Khasi, etc.", "RESOLVED"),
        ("Sensor Hardware Costs", "Budget constraints", "Satellite-first baseline monitoring supplemented by low-cost solar community IoT", "RESOLVED"),
    ]
    for i, (r_name, r_imp, r_mit, r_stat) in enumerate(risks):
        y = Inches(4.45 + i * 0.58)
        add_textbox(s4, Inches(5.35), y, Inches(1.9), Inches(0.25),
                    f"⚠ {r_name}", font_size=8.5, font_color=WHITE, bold=True)
        add_textbox(s4, Inches(7.3), y, Inches(1.4), Inches(0.25),
                    r_imp, font_size=8, font_color=LIGHT_GRAY)
        add_textbox(s4, Inches(8.75), y, Inches(3.1), Inches(0.5),
                    f"✔ {r_mit}", font_size=8, font_color=ACCENT_CYAN)
        add_textbox(s4, Inches(11.9), y, Inches(0.85), Inches(0.25),
                    r_stat, font_size=8, font_color=ACCENT_GREEN, bold=True)

    add_footer_bar(s4, 4, 6)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 5: IMPACT & BENEFITS
    # ──────────────────────────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5, DARK_NAVY)
    add_slide_header(s5, "IMPACT & BENEFITS", "SMART INDIA HACKATHON 2026  |  SOCIO-ECONOMIC IMPACT & SDG ALIGNMENT", ACCENT_ORANGE)

    # 4 Quantifiable Impact Metric Cards
    metrics = [
        ("60%+", "Casualty Reduction", "2 to 6-hour proactive warning window enables safe mass evacuation before collapse.", ACCENT_GREEN),
        ("₹500 Cr+", "Annual Savings", "Mitigates highway transit disruption, equipment loss & post-disaster rescue bills.", ACCENT_BLUE),
        ("45M+", "Citizens Protected", "Comprehensive hazard shielding across all 8 North Eastern states and border villages.", ACCENT_ORANGE),
        (">90%", "Prediction Precision", "Multi-source AI ensemble eliminates false alarms and mitigates alert fatigue.", ACCENT_CYAN),
    ]
    for i, (num, label, desc, col) in enumerate(metrics):
        x = Inches(0.5) + Inches(i * 3.1)
        add_card(s5, x, Inches(1.2), Inches(2.95), Inches(1.85))
        add_shape(s5, MSO_SHAPE.RECTANGLE, x + Inches(0.05), Inches(1.25), Inches(2.85), Pt(3), fill_color=col)
        add_textbox(s5, x + Inches(0.1), Inches(1.35), Inches(2.75), Inches(0.55),
                    num, font_size=28, font_color=col, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s5, x + Inches(0.1), Inches(1.95), Inches(2.75), Inches(0.3),
                    label, font_size=11, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s5, x + Inches(0.1), Inches(2.25), Inches(2.75), Inches(0.7),
                    desc, font_size=8, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # UN Sustainable Development Goals (SDGs)
    add_card(s5, Inches(0.5), Inches(3.2), Inches(12.33), Inches(1.75))
    add_accent_line(s5, Inches(0.6), Inches(3.28), Inches(12.1), color=GOLD)
    add_textbox(s5, Inches(0.7), Inches(3.35), Inches(6), Inches(0.25),
                "ALIGNMENT WITH UN SUSTAINABLE DEVELOPMENT GOALS (SDGs)", font_size=11, font_color=GOLD, bold=True)

    sdgs = [
        ("SDG 11: Sustainable Cities", "Target 11.5: Protect vulnerable hill settlements & reduce disaster casualties.", ACCENT_ORANGE),
        ("SDG 13: Climate Action", "Target 13.1: Adaptive resilience against climate-driven extreme cloudbursts.", ACCENT_GREEN),
        ("SDG 9: Resilient Infrastructure", "Target 9.1: Proactively safeguard highway lifelines (NH-10, NH-29) and bridges.", ACCENT_BLUE),
        ("SDG 17: Multi-Agency Partnerships", "Target 17.16: Unified coordination between MDoNER, NDMA, IMD, GSI & BRO.", ACCENT_CYAN),
    ]
    for i, (s_title, s_desc, s_col) in enumerate(sdgs):
        x = Inches(0.7) + Inches(i * 3.05)
        add_shape(s5, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.68), Inches(2.9), Inches(1.15),
                  fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
        add_shape(s5, MSO_SHAPE.RECTANGLE, x + Inches(0.05), Inches(3.72), Inches(2.8), Pt(2), fill_color=s_col)
        add_textbox(s5, x + Inches(0.1), Inches(3.8), Inches(2.7), Inches(0.35),
                    s_title, font_size=9, font_color=s_col, bold=True)
        add_textbox(s5, x + Inches(0.1), Inches(4.18), Inches(2.7), Inches(0.6),
                    s_desc, font_size=8, font_color=LIGHT_GRAY)

    # Multi-Stakeholder Value Proposition
    add_card(s5, Inches(0.5), Inches(5.1), Inches(12.33), Inches(1.8))
    add_accent_line(s5, Inches(0.6), Inches(5.18), Inches(12.1), color=ACCENT_CYAN)
    add_textbox(s5, Inches(0.7), Inches(5.25), Inches(6), Inches(0.25),
                "MULTI-STAKEHOLDER BENEFICIARIES & VALUE PROPOSITION", font_size=11, font_color=ACCENT_CYAN, bold=True)

    beneficiaries = [
        ("District Admins & SDMAs", "Automated hazard GIS dashboard, resource dispatch optimization, evacuated corridor routing.", ACCENT_BLUE),
        ("Hill Citizens & Farmers", "Life-saving vernacular voice/SMS warnings, livestock/dwelling security, citizen reporting.", ACCENT_GREEN),
        ("MDoNER & Policy Makers", "Granular spatial analytics for targeted slope mitigation funding & cross-state coordination.", ACCENT_ORANGE),
        ("BRO & Transport Lifelines", "Continuous highway vulnerability heatmaps, pre-positioned clearing machinery at high-risk turns.", ACCENT_CYAN),
    ]
    for i, (b_name, b_desc, b_col) in enumerate(beneficiaries):
        x = Inches(0.7) + Inches(i * 3.05)
        add_shape(s5, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.55), Inches(2.9), Inches(1.2),
                  fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
        add_textbox(s5, x + Inches(0.1), Inches(5.62), Inches(2.7), Inches(0.3),
                    b_name, font_size=9.5, font_color=b_col, bold=True)
        add_textbox(s5, x + Inches(0.1), Inches(5.95), Inches(2.7), Inches(0.75),
                    b_desc, font_size=8, font_color=LIGHT_GRAY)

    add_footer_bar(s5, 5, 6)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 6: RESEARCH & REFERENCES
    # ──────────────────────────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6, DARK_NAVY)
    add_slide_header(s6, "RESEARCH, REFERENCES & COMPLIANCE", "SMART INDIA HACKATHON 2026  |  SCIENTIFIC CITATIONS & AUDIT", ACCENT_CYAN)

    # Primary Data Sources
    add_card(s6, Inches(0.5), Inches(1.2), Inches(6.0), Inches(2.6))
    add_accent_line(s6, Inches(0.6), Inches(1.3), Inches(5.8), color=ACCENT_BLUE)
    add_textbox(s6, Inches(0.7), Inches(1.4), Inches(5.6), Inches(0.25),
                "PRIMARY DATA SOURCES & SATELLITE REPOSITORIES", font_size=11, font_color=ACCENT_BLUE, bold=True)

    sources = [
        ("NRSC / ISRO (Bhuvan Portal):", "National Landslide Susceptibility Mapping & Historical Catalog"),
        ("IMD / MoES (Mausam APIs):", "Real-time AWS weather stations, Doppler Radar & precipitation forecasts"),
        ("Geological Survey of India (GSI):", "1:50,000 scale geological hazard & lithological fault line layers"),
        ("Copernicus Open Access (ESA):", "Sentinel-2 Multispectral & Sentinel-1 SAR surface displacement feeds"),
        ("USGS Earth Explorer:", "SRTM 30m Digital Elevation Models (DEM) & Landsat-8/9 thermal data"),
    ]
    for i, (k, v) in enumerate(sources):
        y = Inches(1.72 + i * 0.4)
        add_textbox(s6, Inches(0.7), y, Inches(5.6), Inches(0.38),
                    f"• {k} {v}", font_size=8, font_color=LIGHT_GRAY)

    # Scientific Benchmarks & Literature
    add_card(s6, Inches(6.83), Inches(1.2), Inches(6.0), Inches(2.6))
    add_accent_line(s6, Inches(6.93), Inches(1.3), Inches(5.8), color=ACCENT_GREEN)
    add_textbox(s6, Inches(7.03), Inches(1.4), Inches(5.6), Inches(0.25),
                "KEY SCIENTIFIC BENCHMARKS & LITERATURE", font_size=11, font_color=ACCENT_GREEN, bold=True)

    benchmarks = [
        ("IIT Mandi AI Landslide Network:", "Indigenous MEMS-based early warning deployed in Himalayas (>90% recall)."),
        ("Amrita University IoT LEWS:", "Pioneering multi-tier wireless sensor networks in Sikkim & Western Ghats."),
        ("NIT Silchar Earthquake Studies:", "Empirical studies on seismic-induced slope failure thresholds in Assam/NER."),
        ("MeitY IndiaAI Terralux:", "National open-access AI-powered geospatial platform for disaster management."),
        ("NDMA Guidelines (2019):", "National Landslide Risk Management Strategy for early warning and resilience."),
    ]
    for i, (k, v) in enumerate(benchmarks):
        y = Inches(1.72 + i * 0.4)
        add_textbox(s6, Inches(7.03), y, Inches(5.6), Inches(0.38),
                    f"• {k} {v}", font_size=8, font_color=LIGHT_GRAY)

    # SIH Compliance Checklist (Audit Verification)
    add_card(s6, Inches(0.5), Inches(3.95), Inches(12.33), Inches(1.5))
    add_accent_line(s6, Inches(0.6), Inches(4.03), Inches(12.1), color=GOLD)
    add_textbox(s6, Inches(0.7), Inches(4.1), Inches(6), Inches(0.25),
                "SIH 2026 OFFICIAL SUBMISSION COMPLIANCE CHECKLIST", font_size=11, font_color=GOLD, bold=True)

    compliance = [
        ("Strict 6-Slide Structure", "Exactly 6 slides matching prescribed AICTE / MIC SIH format with zero omissions.", ACCENT_GREEN),
        ("6-Member Team Mandate", "6 members total with female representation (Prerana Mondal) - 100% compliant.", ACCENT_BLUE),
        ("Problem Statement Mapping", "Accurately mapped to SH2001 (Ministry of Development of NER - Software).", ACCENT_CYAN),
        ("Technical & Regional Depth", "NER geology context, AI models, GIS stack, offline PWA & vernacular alerts.", ACCENT_ORANGE),
    ]
    for i, (c_title, c_desc, c_col) in enumerate(compliance):
        x = Inches(0.7) + Inches(i * 3.05)
        add_shape(s6, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.4), Inches(2.9), Inches(0.9),
                  fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
        add_textbox(s6, x + Inches(0.1), Inches(4.46), Inches(2.7), Inches(0.25),
                    f"✔ {c_title}", font_size=9, font_color=c_col, bold=True)
        add_textbox(s6, x + Inches(0.1), Inches(4.72), Inches(2.7), Inches(0.52),
                    c_desc, font_size=7.5, font_color=LIGHT_GRAY)

    # Thank You & Closing Callout
    callout = add_shape(s6, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(5.6), Inches(9.33), Inches(1.25),
                        fill_color=RGBColor(0x0D, 0x47, 0xA1), line_color=ACCENT_CYAN, line_width=2)
    add_multi_text(s6, Inches(2.0), Inches(5.7), Inches(9.33), Inches(1.05), [
        {'text': '🙏 Thank You — Team AlertNex', 'font_size': 20, 'font_color': WHITE,
         'bold': True, 'alignment': PP_ALIGN.CENTER, 'space_after': 4},
        {'text': '"Empowering NER with AI-driven disaster resilience for a safer tomorrow"',
         'font_size': 12, 'font_color': ACCENT_CYAN, 'italic': True, 'alignment': PP_ALIGN.CENTER, 'space_after': 2},
        {'text': 'Smart India Hackathon 2026  |  Ministry of Development of North Eastern Region (MDoNER)',
         'font_size': 10, 'font_color': LIGHT_GRAY, 'alignment': PP_ALIGN.CENTER},
    ])

    add_footer_bar(s6, 6, 6)

    # Save
    output_dir = os.path.dirname(os.path.abspath(__file__))
    p1 = os.path.join(output_dir, "SIH2026_AlertNex_Presentation.pptx")
    prs.save(p1)
    p2 = os.path.join(output_dir, "SIH2026_AlertNex_Official_6Slides.pptx")
    prs.save(p2)
    print("Official 6-Slide Presentation saved successfully:", p1)
    return p1


# ==============================================================================
# 2. EXTENDED 10-SLIDE PRESENTATION BUILDER (FOR INTERNAL / EXTENDED ROUNDS)
# ==============================================================================
def create_extended_10slide_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # Slide 1: Title
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1, DARK_NAVY)
    add_shape(s1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), fill_color=ACCENT_CYAN)

    b1 = add_shape(s1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.4), Inches(3.2), Inches(0.45),
                   fill_color=RGBColor(0x0D, 0x47, 0xA1), line_color=ACCENT_BLUE, line_width=1)
    b1.text_frame.paragraphs[0].text = "SMART INDIA HACKATHON 2026"
    b1.text_frame.paragraphs[0].font.size = Pt(12)
    b1.text_frame.paragraphs[0].font.color.rgb = WHITE
    b1.text_frame.paragraphs[0].font.bold = True
    b1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    b2 = add_shape(s1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.85), Inches(0.4), Inches(1.8), Inches(0.45), fill_color=ACCENT_ORANGE)
    b2.text_frame.paragraphs[0].text = "PS: SH2001"
    b2.text_frame.paragraphs[0].font.size = Pt(12)
    b2.text_frame.paragraphs[0].font.color.rgb = WHITE
    b2.text_frame.paragraphs[0].font.bold = True
    b2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_multi_text(s1, Inches(0.5), Inches(1.1), Inches(8.5), Inches(2.2), [
        {'text': 'AI-Based Early Warning &', 'font_size': 34, 'font_color': WHITE, 'bold': True, 'space_after': 4},
        {'text': 'Landslide Risk Monitoring System', 'font_size': 34, 'font_color': ACCENT_CYAN, 'bold': True, 'space_after': 4},
        {'text': 'for North Eastern Region (NER) of India', 'font_size': 20, 'font_color': LIGHT_GRAY, 'space_after': 6},
    ])

    info_card = add_card(s1, Inches(9.2), Inches(0.95), Inches(3.6), Inches(2.3))
    add_accent_line(s1, Inches(9.3), Inches(1.05), Inches(3.4), color=ACCENT_CYAN)
    add_textbox(s1, Inches(9.4), Inches(1.15), Inches(3.2), Inches(0.3),
                "OFFICIAL PROBLEM DETAILS", font_size=11, font_color=ACCENT_CYAN, bold=True)
    ps_info = [
        ("Organization:", "Ministry of Development of NER"),
        ("Problem Statement:", "SH2001 (Software Edition)"),
        ("Theme:", "Disaster Management"),
        ("Institute Name:", "[Your College / Institute Name, City]"),
        ("Faculty Mentor:", "[Faculty Mentor Name / Designation]"),
    ]
    for i, (k, v) in enumerate(ps_info):
        add_textbox(s1, Inches(9.4), Inches(1.5 + i * 0.32), Inches(1.3), Inches(0.3), k, font_size=9, font_color=MEDIUM_GRAY)
        add_textbox(s1, Inches(10.6), Inches(1.5 + i * 0.32), Inches(2.1), Inches(0.3), v, font_size=9, font_color=WHITE, bold=True)

    add_accent_line(s1, Inches(0.5), Inches(3.3), Inches(4.5), color=ACCENT_ORANGE)
    add_textbox(s1, Inches(0.5), Inches(3.4), Inches(5.5), Inches(0.4),
                "Team AlertNex  (Lead: Ayush Kumar)", font_size=20, font_color=ACCENT_ORANGE, bold=True)
    add_textbox(s1, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.35),
                "\"Turning raw environmental data into actionable intelligence for landslide disaster resilience in NER\"",
                font_size=12, font_color=ACCENT_GREEN, italic=True)

    # 6 Members Card
    add_card(s1, Inches(0.5), Inches(4.3), Inches(12.33), Inches(2.6))
    add_accent_line(s1, Inches(0.6), Inches(4.38), Inches(12.1), color=ACCENT_CYAN)
    add_textbox(s1, Inches(0.7), Inches(4.45), Inches(6), Inches(0.3),
                "TEAM MEMBERS  (SIH 6-Member Mandate Compliant)", font_size=11, font_color=ACCENT_CYAN, bold=True)
    add_textbox(s1, Inches(7.0), Inches(4.45), Inches(5.6), Inches(0.3),
                "✔ Female Representation: Prerana Mondal  |  100% SIH Compliant", font_size=10, font_color=GOLD, italic=True, alignment=PP_ALIGN.RIGHT)

    for i, (name, role, skills, color) in enumerate(TEAM_MEMBERS):
        x = Inches(0.65) + Inches(i * 2.02)
        add_shape(s1, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.85), Inches(1.95), Inches(1.95),
                  fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
        add_shape(s1, MSO_SHAPE.RECTANGLE, x + Inches(0.05), Inches(4.88), Inches(1.85), Pt(2), fill_color=color)

        av = add_shape(s1, MSO_SHAPE.OVAL, x + Inches(0.65), Inches(4.98), Inches(0.65), Inches(0.65), fill_color=color)
        av.text_frame.paragraphs[0].text = name[2] if name[2] != '[' else '6'
        av.text_frame.paragraphs[0].font.size = Pt(16)
        av.text_frame.paragraphs[0].font.color.rgb = WHITE
        av.text_frame.paragraphs[0].font.bold = True
        av.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(s1, x + Inches(0.05), Inches(5.7), Inches(1.85), Inches(0.3),
                    name.replace("👤 ", ""), font_size=10, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s1, x + Inches(0.05), Inches(5.98), Inches(1.85), Inches(0.25),
                    role, font_size=8.5, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s1, x + Inches(0.05), Inches(6.22), Inches(1.85), Inches(0.55),
                    skills, font_size=7.5, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_footer_bar(s1, 1, 10)

    # Slide 2: Problem Understanding
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2, DARK_NAVY)
    add_slide_header(s2, "PROBLEM UNDERSTANDING & BACKGROUND", "SMART INDIA HACKATHON 2026", ACCENT_ORANGE)

    stats = [
        ("4,796+", "Landslide Events", "recorded in NER\n(2015-2024)", ACCENT_RED),
        ("~1,200", "Deaths Annually", "due to landslides\nacross India", ACCENT_ORANGE),
        ("8 States", "Affected Region", "Arunachal, Assam,\nManipur, Meghalaya,\nMizoram, Nagaland,\nSikkim, Tripura", ACCENT_BLUE),
        ("70%+", "Reactive Response", "Current systems rely\non manual reporting", ACCENT_CYAN),
    ]
    for i, (num, title, desc, col) in enumerate(stats):
        x = Inches(0.5) + Inches(i * 3.1)
        add_card(s2, x, Inches(1.2), Inches(2.8), Inches(2.2))
        add_accent_line(s2, x + Inches(0.05), Inches(1.25), Inches(2.7), color=col)
        add_textbox(s2, x + Inches(0.2), Inches(1.4), Inches(2.4), Inches(0.5), num, font_size=32, font_color=col, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s2, x + Inches(0.2), Inches(1.95), Inches(2.4), Inches(0.35), title, font_size=13, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s2, x + Inches(0.2), Inches(2.35), Inches(2.4), Inches(0.8), desc, font_size=9.5, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(s2, Inches(0.5), Inches(3.65), Inches(12), Inches(0.35), "THE CORE PROBLEM STATEMENT IN NER", font_size=15, font_color=ACCENT_CYAN, bold=True)
    problems = [
        "• Fragile Himalayan Geology: Young, seismically active terrain with steep slopes prone to catastrophic debris flows and flash floods.",
        "• Monsoonal Downpours: Over 2,000 mm annual rainfall triggers rapid soil saturation and elevated pore pressure in hours.",
        "• Highway Severance: Arterial lifelines NH-10 (Sikkim) and NH-29 (Nagaland/Manipur) repeatedly blocked, crippling supply chains.",
        "• Reactive Response: Disaster management currently depends on post-event manual reporting with no unified predictive platform.",
        "• Severe Telecom Barriers: Hilly terrain suffers frequent cellular outages during storms, halting warnings without offline failover.",
    ]
    for i, p in enumerate(problems):
        add_textbox(s2, Inches(0.5), Inches(4.15 + i * 0.48), Inches(12.3), Inches(0.45), p, font_size=11, font_color=LIGHT_GRAY)
    add_footer_bar(s2, 2, 10)

    # Slide 3: Proposed Solution
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3, DARK_NAVY)
    add_slide_header(s3, "PROPOSED SOLUTION — AlertNex PLATFORM", "SMART INDIA HACKATHON 2026", ACCENT_GREEN)

    sol_cards = [
        ("📡", "Multi-Source Data Ingestion", "Automated pipelines fusing IMD radar, Sentinel-2 optical/SAR imagery, GSI lithology, and IoT soil readings.", ACCENT_BLUE),
        ("🧠", "Hybrid AI/ML Predictive Engine", "Ensemble of Random Forest, XGBoost, CNN displacement, and LSTM forecasting slope failure (>90% accuracy).", ACCENT_CYAN),
        ("🗺️", "Dynamic GIS Risk Heatmaps", "Real-time interactive Leaflet & PostGIS visualizer mapping live hazard danger zones & evacuation routes.", ACCENT_GREEN),
        ("🔔", "Multi-Channel Local Alerts", "Disseminates warnings via SMS, siren triggers, and synthesized voice calls in native NER dialects.", ACCENT_ORANGE),
        ("📱", "Citizen Crowdsourcing PWA", "Progressive Web App empowering citizens to submit geo-tagged observations of ground cracks & seepage.", GOLD),
        ("📶", "Offline & Low-Network Mode", "Edge caching, store-and-forward sync, and SMS fallback ensuring 100% operation during telecom outages.", ACCENT_RED),
    ]
    for i, (icon, title, desc, col) in enumerate(sol_cards):
        col_idx = i % 3
        row_idx = i // 3
        x = Inches(0.5) + Inches(col_idx * 4.2)
        y = Inches(1.3) + Inches(row_idx * 2.7)
        add_card(s3, x, y, Inches(3.9), Inches(2.4))
        add_accent_line(s3, x + Inches(0.1), y + Inches(0.05), Inches(3.7), color=col)

        ic = add_shape(s3, MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.2), Inches(0.65), Inches(0.65), fill_color=col)
        ic.text_frame.paragraphs[0].text = icon
        ic.text_frame.paragraphs[0].font.size = Pt(18)
        ic.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(s3, x + Inches(0.2), y + Inches(0.95), Inches(3.5), Inches(0.4), title, font_size=13, font_color=col, bold=True)
        add_textbox(s3, x + Inches(0.2), y + Inches(1.35), Inches(3.5), Inches(0.95), desc, font_size=10, font_color=LIGHT_GRAY)
    add_footer_bar(s3, 3, 10)

    # Slide 4: Technical Architecture
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4, DARK_NAVY)
    add_slide_header(s4, "TECHNICAL ARCHITECTURE & DATA FLOW", "SMART INDIA HACKATHON 2026", ACCENT_BLUE)

    layers_ext = [
        ("Layer 1: Ingestion", ACCENT_BLUE, ["Sentinel-2 & Landsat Satellites", "IMD Rainfall Radar APIs", "GSI Lithology & SRTM DEM", "Citizen Geo-tagged Reports"]),
        ("Layer 2: AI Analytics", ACCENT_CYAN, ["Feature Extraction (NDVI, Slope)", "Random Forest & XGBoost Ensemble", "CNN Terrain Displacement Model", "LSTM Rainfall Trigger Forecasting"]),
        ("Layer 3: Geospatial", ACCENT_GREEN, ["PostgreSQL / PostGIS Database", "GeoServer Vector Map Tiles", "Shortest Safe Evacuation Routing", "Multilingual Voice & NLP Engine"]),
        ("Layer 4: User Apps", ACCENT_ORANGE, ["Admin Command Dashboard (React)", "Citizen & Field PWA (Offline Sync)", "Public Travel Advisory Portal", "REST APIs for NDMA / BRO"]),
    ]
    for i, (l_name, col, items) in enumerate(layers_ext):
        x = Inches(0.5) + Inches(i * 3.1)
        add_card(s4, x, Inches(1.3), Inches(2.95), Inches(5.4))
        add_accent_line(s4, x + Inches(0.05), Inches(1.35), Inches(2.85), color=col)
        add_textbox(s4, x + Inches(0.1), Inches(1.5), Inches(2.75), Inches(0.4), l_name, font_size=12, font_color=col, bold=True, alignment=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_shape(s4, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), Inches(2.2 + j * 1.1), Inches(2.65), Inches(0.9),
                      fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
            add_textbox(s4, x + Inches(0.2), Inches(2.3 + j * 1.1), Inches(2.55), Inches(0.7), item, font_size=9, font_color=WHITE, alignment=PP_ALIGN.CENTER)
    add_footer_bar(s4, 4, 10)

    # Slide 5: Innovation & Novelty
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5, DARK_NAVY)
    add_slide_header(s5, "INNOVATION, NOVELTY & COMPETITIVE EDGE", "SMART INDIA HACKATHON 2026", GOLD)

    innovations = [
        ("Multi-Modal AI Ensemble", "Fuses optical satellite imagery, radar InSAR, and ground sensors into an ensemble model achieving >90% precision with low false alarms.", ACCENT_BLUE),
        ("Offline-First Resilience", "Local edge caching with IndexedDB and store-and-forward sync ensures alerts and reporting operate without cellular data.", ACCENT_CYAN),
        ("Vernacular Dialect Alerts", "Automated text and voice synthesis in 10+ North Eastern languages (Assamese, Mizo, Khasi, Manipuri) directly reaching indigenous communities.", ACCENT_GREEN),
        ("Closed-Loop Crowdsourcing", "Citizen reports of ground fissures and seepage feed back into models to continuously calibrate and fine-tune hazard boundaries.", ACCENT_ORANGE),
    ]
    for i, (title, desc, col) in enumerate(innovations):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(0.5) + Inches(col_idx * 6.3)
        y = Inches(1.4) + Inches(row_idx * 2.6)
        add_card(s5, x, y, Inches(6.0), Inches(2.3))
        add_accent_line(s5, x + Inches(0.1), y + Inches(0.05), Inches(5.8), color=col)
        add_textbox(s5, x + Inches(0.3), y + Inches(0.3), Inches(5.4), Inches(0.4), title, font_size=15, font_color=col, bold=True)
        add_textbox(s5, x + Inches(0.3), y + Inches(0.8), Inches(5.4), Inches(1.2), desc, font_size=11, font_color=LIGHT_GRAY)
    add_footer_bar(s5, 5, 10)

    # Slide 6: Feasibility & Viability
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6, DARK_NAVY)
    add_slide_header(s6, "FEASIBILITY & VIABILITY MATRIX", "SMART INDIA HACKATHON 2026", ACCENT_GREEN)

    add_card(s6, Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.4))
    add_accent_line(s6, Inches(0.6), Inches(1.35), Inches(5.7), color=ACCENT_GREEN)
    add_textbox(s6, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.35), "HACKATHON FEASIBILITY HIGHLIGHTS", font_size=13, font_color=ACCENT_GREEN, bold=True)

    feas_points = [
        ("Pre-trained Open Models:", "Using proven XGBoost & Random Forest baselines trained on NRSC catalogs eliminates cold-start training delays."),
        ("Open-Source Geospatial:", "Leaflet + PostGIS + GeoServer enables rapid implementation without proprietary software licensing fees."),
        ("Cloud Containerization:", "Dockerized microservices allow turnkey deployment across testbeds in minutes."),
        ("Zero Licensing Costs:", "100% open-source software stack ensures unmatched economic viability for government scaling."),
    ]
    for i, (k, v) in enumerate(feas_points):
        y = Inches(2.0 + i * 1.1)
        add_textbox(s6, Inches(0.7), y, Inches(5.5), Inches(0.3), f"✔ {k}", font_size=10.5, font_color=WHITE, bold=True)
        add_textbox(s6, Inches(0.85), y + Inches(0.28), Inches(5.35), Inches(0.7), v, font_size=9.5, font_color=LIGHT_GRAY)

    add_card(s6, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.4))
    add_accent_line(s6, Inches(6.9), Inches(1.35), Inches(5.8), color=ACCENT_RED)
    add_textbox(s6, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.35), "CHALLENGES & ENGINEERED MITIGATIONS", font_size=13, font_color=ACCENT_RED, bold=True)

    challenges = [
        ("Hill Telecom Blackouts", "Offline-first PWA caching + store-and-forward sync + automated SMS broadcast fallback."),
        ("Complex Micro-Climates", "Ensemble AI fusing local rain gauges with satellite InSAR and DEM slope analysis."),
        ("Ethnic Linguistic Diversity", "Automated text & speech synthesis in 10+ local North Eastern dialects."),
        ("Sensor Hardware Expenses", "Satellite-first baseline monitoring supplemented by low-cost community solar sensors."),
    ]
    for i, (ch, mit) in enumerate(challenges):
        y = Inches(2.0 + i * 1.1)
        add_textbox(s6, Inches(7.0), y, Inches(5.6), Inches(0.3), f"⚠ {ch}", font_size=10.5, font_color=ACCENT_ORANGE, bold=True)
        add_textbox(s6, Inches(7.15), y + Inches(0.28), Inches(5.45), Inches(0.7), f"Mitigation: {mit}", font_size=9.5, font_color=LIGHT_GRAY)
    add_footer_bar(s6, 6, 10)

    # Slide 7: Impact & Benefits
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7, DARK_NAVY)
    add_slide_header(s7, "IMPACT, BENEFITS & UN SDG ALIGNMENT", "SMART INDIA HACKATHON 2026", ACCENT_CYAN)

    for i, (num, label, desc, col) in enumerate(METRICS):
        x = Inches(0.5) + Inches(i * 3.1)
        add_card(s7, x, Inches(1.3), Inches(2.95), Inches(2.2))
        add_accent_line(s7, x + Inches(0.05), Inches(1.35), Inches(2.85), color=col)
        add_textbox(s7, x + Inches(0.1), Inches(1.5), Inches(2.75), Inches(0.6), num, font_size=32, font_color=col, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s7, x + Inches(0.1), Inches(2.2), Inches(2.75), Inches(0.35), label, font_size=12, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s7, x + Inches(0.1), Inches(2.6), Inches(2.75), Inches(0.8), desc, font_size=8.5, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_card(s7, Inches(0.5), Inches(3.8), Inches(12.33), Inches(2.9))
    add_accent_line(s7, Inches(0.6), Inches(3.88), Inches(12.1), color=GOLD)
    add_textbox(s7, Inches(0.7), Inches(3.98), Inches(6), Inches(0.3), "ALIGNMENT WITH UN SUSTAINABLE DEVELOPMENT GOALS", font_size=12, font_color=GOLD, bold=True)

    for i, (s_title, s_desc, s_col) in enumerate(SDGS):
        x = Inches(0.7) + Inches(i * 3.05)
        add_shape(s7, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.4), Inches(2.9), Inches(2.0),
                  fill_color=RGBColor(0x13, 0x33, 0x55), line_color=CARD_BORDER, line_width=1)
        add_shape(s7, MSO_SHAPE.RECTANGLE, x + Inches(0.05), Inches(4.45), Inches(2.8), Pt(2), fill_color=s_col)
        add_textbox(s7, x + Inches(0.15), Inches(4.6), Inches(2.6), Inches(0.45), s_title, font_size=10.5, font_color=s_col, bold=True)
        add_textbox(s7, x + Inches(0.15), Inches(5.15), Inches(2.6), Inches(1.1), s_desc, font_size=8.5, font_color=LIGHT_GRAY)
    add_footer_bar(s7, 7, 10)

    # Slide 8: Prototype UI/UX Wireframes
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8, DARK_NAVY)
    add_slide_header(s8, "PROTOTYPE ARCHITECTURE & UI/UX WIREFRAMES", "SMART INDIA HACKATHON 2026", ACCENT_BLUE)

    ui_modules = [
        ("Admin Command Console", "React.js dashboard displaying live interactive Leaflet risk heatmap, weather overlays, and emergency resource dispatch panels.", ACCENT_BLUE),
        ("Citizen Mobile PWA", "Offline-first application for local hill residents with geo-tagged crack reporting, multilingual voice alerts, and safe evacuation maps.", ACCENT_GREEN),
        ("BRO & Transport Portal", "Highway vulnerability tracking for NH-10 and NH-29, showing clearance machinery pre-positioning and live road status.", ACCENT_ORANGE),
    ]
    for i, (u_title, u_desc, u_col) in enumerate(ui_modules):
        x = Inches(0.5) + Inches(i * 4.2)
        add_card(s8, x, Inches(1.4), Inches(3.9), Inches(5.3))
        add_accent_line(s8, x + Inches(0.1), Inches(1.45), Inches(3.7), color=u_col)
        add_textbox(s8, x + Inches(0.2), Inches(1.6), Inches(3.5), Inches(0.35), u_title, font_size=13, font_color=u_col, bold=True)
        add_textbox(s8, x + Inches(0.2), Inches(2.05), Inches(3.5), Inches(0.75), u_desc, font_size=9.5, font_color=LIGHT_GRAY)

        mock = add_shape(s8, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), Inches(2.9), Inches(3.5), Inches(3.4),
                         fill_color=RGBColor(0x08, 0x18, 0x2A), line_color=u_col, line_width=1)
        add_textbox(s8, x + Inches(0.3), Inches(3.2), Inches(3.3), Inches(0.3), "MOCKUP PREVIEW", font_size=10, font_color=u_col, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s8, x + Inches(0.3), Inches(3.7), Inches(3.3), Inches(2.2),
                    f"• Interactive GIS layers\n• Real-time risk gauge\n• SMS Broadcast trigger\n• Offline cache indicator\n• Evacuation routing",
                    font_size=9, font_color=LIGHT_GRAY)
    add_footer_bar(s8, 8, 10)

    # Slide 9: 36-Hour Development Timeline
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9, DARK_NAVY)
    add_slide_header(s9, "36-HOUR HACKATHON DEVELOPMENT TIMELINE", "SMART INDIA HACKATHON 2026", ACCENT_GREEN)

    timeline_ext = [
        ("Phase 1", "0-8 Hours", "Setup & ML Modeling", ["Ingest NRSC historical landslide inventory", "Pre-train Random Forest & XGBoost baseline", "Establish Git CI/CD & Docker environment"], ACCENT_BLUE),
        ("Phase 2", "8-18 Hours", "GIS & Spatial Engine", ["Configure PostgreSQL / PostGIS spatial DB", "Deploy GeoServer & MapLibre raster tiles", "Compute DEM slope & aspect algorithms"], ACCENT_CYAN),
        ("Phase 3", "18-28 Hours", "UI & Alert Gateways", ["Build React Command Console & PWA UI", "Integrate Twilio SMS & FCM push alerts", "Implement offline service worker cache"], ACCENT_GREEN),
        ("Phase 4", "28-36 Hours", "Validation & Pitch", ["Simulate Sikkim cloudburst disaster scenario", "Stress-test offline failover under zero cellular", "Finalize live pitch demo & documentation"], ACCENT_ORANGE),
    ]
    for i, (ph, hrs, t_title, tasks, col) in enumerate(timeline_ext):
        y = Inches(1.4 + i * 1.35)
        add_card(s9, Inches(0.5), y, Inches(12.33), Inches(1.2))
        add_accent_line(s9, Inches(0.6), y + Inches(0.05), Inches(12.1), color=col)

        b = add_shape(s9, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y + Inches(0.15), Inches(1.3), Inches(0.35), fill_color=col)
        b.text_frame.paragraphs[0].text = ph
        b.text_frame.paragraphs[0].font.size = Pt(10)
        b.text_frame.paragraphs[0].font.color.rgb = WHITE
        b.text_frame.paragraphs[0].font.bold = True
        b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(s9, Inches(2.2), y + Inches(0.15), Inches(1.5), Inches(0.3), hrs, font_size=11, font_color=col, bold=True)
        add_textbox(s9, Inches(3.8), y + Inches(0.15), Inches(4.0), Inches(0.3), t_title, font_size=12, font_color=WHITE, bold=True)

        for j, task in enumerate(tasks):
            add_textbox(s9, Inches(0.8) + Inches(j * 3.9), y + Inches(0.55), Inches(3.8), Inches(0.5),
                        f"▸ {task}", font_size=8.5, font_color=LIGHT_GRAY)
    add_footer_bar(s9, 9, 10)

    # Slide 10: Team, References & Closing
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10, DARK_NAVY)
    add_slide_header(s10, "TEAM AlertNex, REFERENCES & SIH AUDIT", "SMART INDIA HACKATHON 2026", ACCENT_CYAN)

    for i, (name, role, skills, color) in enumerate(TEAM_MEMBERS):
        x = Inches(0.5) + Inches(i * 2.08)
        add_card(s10, x, Inches(1.2), Inches(1.98), Inches(2.2))
        add_accent_line(s10, x + Inches(0.05), Inches(1.25), Inches(1.88), color=color)

        av = add_shape(s10, MSO_SHAPE.OVAL, x + Inches(0.65), Inches(1.35), Inches(0.65), Inches(0.65), fill_color=color)
        av.text_frame.paragraphs[0].text = name[2] if name[2] != '[' else '6'
        av.text_frame.paragraphs[0].font.size = Pt(16)
        av.text_frame.paragraphs[0].font.color.rgb = WHITE
        av.text_frame.paragraphs[0].font.bold = True
        av.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(s10, x + Inches(0.05), Inches(2.05), Inches(1.88), Inches(0.25),
                    name.replace("👤 ", ""), font_size=9.5, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s10, x + Inches(0.05), Inches(2.3), Inches(1.88), Inches(0.25),
                    role, font_size=8, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s10, x + Inches(0.05), Inches(2.55), Inches(1.88), Inches(0.75),
                    skills, font_size=7, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # References bottom
    add_card(s10, Inches(0.5), Inches(3.65), Inches(12.33), Inches(1.8))
    add_accent_line(s10, Inches(0.6), Inches(3.72), Inches(12.1), color=ACCENT_CYAN)
    add_textbox(s10, Inches(0.7), Inches(3.8), Inches(6), Inches(0.25), "KEY REFERENCES & SATELLITE DATA SOURCES", font_size=11, font_color=ACCENT_CYAN, bold=True)

    r_l = [
        "• NRSC / ISRO — Bhuvan Landslide Inventory Database",
        "• IMD / MoES — Automatic Weather Station & Doppler Radar APIs",
        "• Geological Survey of India (GSI) — 1:50,000 Susceptibility Maps",
    ]
    r_r = [
        "• IIT Mandi & Amrita University — AI/IoT Landslide EWS Field Research",
        "• NDMA Guidelines — National Landslide Risk Management Strategy",
        "• Copernicus Open Access Hub — Sentinel-1/2 SAR & Optical Satellites",
    ]
    for i, r in enumerate(r_l):
        add_textbox(s10, Inches(0.7), Inches(4.15 + i * 0.35), Inches(5.8), Inches(0.3), r, font_size=8.5, font_color=LIGHT_GRAY)
    for i, r in enumerate(r_r):
        add_textbox(s10, Inches(6.8), Inches(4.15 + i * 0.35), Inches(5.8), Inches(0.3), r, font_size=8.5, font_color=LIGHT_GRAY)

    # Closing Box
    add_shape(s10, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(5.65), Inches(9.33), Inches(1.2),
              fill_color=RGBColor(0x0D, 0x47, 0xA1), line_color=ACCENT_CYAN, line_width=2)
    add_multi_text(s10, Inches(2.0), Inches(5.75), Inches(9.33), Inches(1.0), [
        {'text': '🙏 Thank You — Team AlertNex', 'font_size': 20, 'font_color': WHITE,
         'bold': True, 'alignment': PP_ALIGN.CENTER, 'space_after': 4},
        {'text': '"Empowering NER with AI-driven disaster resilience for a safer tomorrow"',
         'font_size': 12, 'font_color': ACCENT_CYAN, 'italic': True, 'alignment': PP_ALIGN.CENTER},
    ])

    add_footer_bar(s10, 10, 10)

    output_dir = os.path.dirname(os.path.abspath(__file__))
    p_ext = os.path.join(output_dir, "SIH2026_AlertNex_Extended_10Slides.pptx")
    prs.save(p_ext)
    print("Extended 10-Slide Presentation saved successfully:", p_ext)
    return p_ext


def main():
    create_official_6slide_presentation()
    create_extended_10slide_presentation()


if __name__ == "__main__":
    main()

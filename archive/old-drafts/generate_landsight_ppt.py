"""
LANDSIGHT AI - SIH 2026 Official 6-Slide Presentation Generator
Problem Statement ID: SIH26001
Theme: Disaster Management | Category: Software
Organization: Ministry of Development of North Eastern Region (MDoNER)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
BG_DARK = RGBColor(0x0A, 0x19, 0x2F)        # Deep Navy / Space cadet
CARD_BG = RGBColor(0x11, 0x22, 0x40)        # Navy Card Surface
CARD_BORDER = RGBColor(0x23, 0x35, 0x54)    # Card Border
CARD_HOVER = RGBColor(0x1B, 0x3A, 0x60)     # Lighter Card Accent
TEXT_WHITE = RGBColor(0xF8, 0xF9, 0xFA)     # Primary White
TEXT_MUTED = RGBColor(0x88, 0x92, 0xB0)     # Cool Gray
TEXT_ACCENT = RGBColor(0x64, 0xFF, 0xDA)    # Electric Teal / Cyan
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)    # Bright Cyan
ACCENT_BLUE = RGBColor(0x00, 0x99, 0xFF)    # Vivid Blue
ACCENT_GOLD = RGBColor(0xFF, 0xC1, 0x07)    # Gold
RISK_GREEN = RGBColor(0x00, 0xE6, 0x76)     # Low Risk Green
RISK_YELLOW = RGBColor(0xFF, 0xD6, 0x00)    # Moderate Risk Yellow
RISK_ORANGE = RGBColor(0xFF, 0x91, 0x00)    # High Risk Orange
RISK_RED = RGBColor(0xFF, 0x3D, 0x00)       # Critical Risk Red

# Dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

def set_slide_background(slide):
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG_DARK
    bg_shape.line.fill.background()
    return bg_shape

def add_header(slide, title_text, category_text="SIH 2026 | PS ID: SIH26001 | MDoNER"):
    # Header tag badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(4.5), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = CARD_BG
    badge.line.color.rgb = ACCENT_BLUE
    badge.line.width = Pt(1)
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = category_text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Main Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.7))
    tf2 = title_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title_text
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE

    # Bottom footer watermark
    ft = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.35))
    tf_ft = ft.text_frame
    p_ft = tf_ft.paragraphs[0]
    p_ft.text = "LANDSIGHT AI  •  Predict Risk. Protect Lives. Connect Communities.  •  Smart India Hackathon 2026"
    p_ft.font.size = Pt(9.5)
    p_ft.font.color.rgb = TEXT_MUTED

# ==========================================
# SLIDE 1: TITLE SLIDE
# ==========================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1)

# Subtle decorative accent bar at top
top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = ACCENT_BLUE
top_bar.line.fill.background()

# Badges line
s1_badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(0.7), Inches(6.0), Inches(0.42))
s1_badge.fill.solid()
s1_badge.fill.fore_color.rgb = CARD_BG
s1_badge.line.color.rgb = ACCENT_BLUE
s1_badge.line.width = Pt(1)
tf1 = s1_badge.text_frame
tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf1.paragraphs[0]
p.text = "SMART INDIA HACKATHON 2026  |  COLLEGE INTERNAL EVALUATION"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = TEXT_ACCENT
p.alignment = PP_ALIGN.CENTER

# Main Project Name
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.25), Inches(11.33), Inches(1.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "LANDSIGHT AI"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE

p_sub = tf.add_paragraph()
p_sub.text = "AI-Based Early Warning & Landslide Monitoring System for North Eastern Region"
p_sub.font.size = Pt(19)
p_sub.font.bold = True
p_sub.font.color.rgb = ACCENT_BLUE

p_tag = tf.add_paragraph()
p_tag.text = '“Predict Risk. Protect Lives. Connect Communities.”'
p_tag.font.size = Pt(14)
p_tag.font.italic = True
p_tag.font.color.rgb = ACCENT_GOLD

# Meta Details Card (Left)
meta_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.2), Inches(5.4), Inches(3.3))
meta_card.fill.solid()
meta_card.fill.fore_color.rgb = CARD_BG
meta_card.line.color.rgb = CARD_BORDER
meta_card.line.width = Pt(1)
m_tf = meta_card.text_frame
m_tf.word_wrap = True
m_tf.margin_left = Inches(0.25)
m_tf.margin_top = Inches(0.2)

p = m_tf.paragraphs[0]
p.text = "PROJECT SPECIFICATIONS"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = TEXT_ACCENT

specs = [
    ("Problem ID:", "SIH26001"),
    ("Theme:", "Disaster Management"),
    ("Category:", "Software Edition"),
    ("Ministry:", "Ministry of Development of North Eastern Region (MDoNER)"),
    ("Target Focus:", "8 North Eastern States (Sikkim, Assam, Meghalaya, etc.)"),
    ("Key Innovation:", "Dynamic Risk + Connectivity Impact + Offline Sync")
]
for label, val in specs:
    p = m_tf.add_paragraph()
    p.text = f"•  {label} "
    p.font.bold = True
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    run = p.add_run()
    run.text = val
    run.font.bold = False
    run.font.color.rgb = TEXT_MUTED

# Team Details Card (Right)
team_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.2), Inches(5.5), Inches(3.3))
team_card.fill.solid()
team_card.fill.fore_color.rgb = CARD_BG
team_card.line.color.rgb = CARD_BORDER
team_card.line.width = Pt(1)
t_tf = team_card.text_frame
t_tf.word_wrap = True
t_tf.margin_left = Inches(0.25)
t_tf.margin_top = Inches(0.2)

p = t_tf.paragraphs[0]
p.text = "SUBMISSION & TEAM PROFILE"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_GOLD

p = t_tf.add_paragraph()
p.text = "Team Name: [TEAM NAME]  |  Internal Evaluation"
p.font.bold = True
p.font.size = Pt(11)
p.font.color.rgb = TEXT_WHITE

p = t_tf.add_paragraph()
p.text = "Institution: [COLLEGE / INSTITUTION NAME]"
p.font.size = Pt(10.5)
p.font.color.rgb = TEXT_ACCENT

p = t_tf.add_paragraph()
p.text = "Team Roster (6 Members):"
p.font.bold = True
p.font.size = Pt(10.5)
p.font.color.rgb = TEXT_WHITE

team_members = [
    ("1. [Leader Name]", "Team Leader & AI/ML Architecture"),
    ("2. [Member 2 Name]", "Frontend & UI/UX Specialist"),
    ("3. [Member 3 Name]", "Backend API & PostGIS Dev"),
    ("4. [Member 4 Name]", "Mobile App (Flutter) & Offline Storage"),
    ("5. [Member 5 Name]", "Geospatial & Road Network Analyst"),
    ("6. [Member 6 Name]", "Cloud, Edge & Documentation")
]
for name, role in team_members:
    p = t_tf.add_paragraph()
    p.text = f"• {name} — "
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    r = p.add_run()
    r.text = role
    r.font.size = Pt(9.5)
    r.font.color.rgb = TEXT_MUTED

# ==========================================
# SLIDE 2: THE CHALLENGE
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_header(slide2, "THE CHALLENGE: Critical Landslide Realities in North Eastern Region")

# Problem statement summary banner
sum_banner = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.7))
sum_banner.fill.solid()
sum_banner.fill.fore_color.rgb = CARD_BG
sum_banner.line.color.rgb = RISK_RED
sum_banner.line.width = Pt(1)
sb_tf = sum_banner.text_frame
sb_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
sb_tf.margin_left = Inches(0.2)
p = sb_tf.paragraphs[0]
p.text = "PROBLEM SUMMARY: "
p.font.bold = True
p.font.size = Pt(11)
p.font.color.rgb = RISK_RED
r = p.add_run()
r.text = "NER faces extreme rainfall & steep geology, causing catastrophic landslides that sever critical arterial highways (NH-10, NH-29), isolate remote tribal villages, and cut hospital access with zero proactive connectivity intelligence."
r.font.bold = False
r.font.color.rgb = TEXT_WHITE

# 3 Critical Problem Pillars (Cards)
cols = [
    ("1. GEOGRAPHICAL VULNERABILITY", [
        "Heavy monsoon & cloudbursts exceed soil shear strength",
        "Steep terrain & young Himalayas cause sudden slope slips",
        "Over 70% of NER land is landslide-prone mountain zone"
    ], RISK_ORANGE),
    ("2. INFRASTRUCTURE & COMMUNICATION", [
        "Remote villages depend on single lifeline mountain roads",
        "Unstable cellular network leaves zero connectivity in rain",
        "Rockfalls cut power, cellular towers, and medical corridors"
    ], RISK_YELLOW),
    ("3. DECISION SYSTEM GAPS", [
        "Authorities receive ground reports hours/days too late",
        "Existing systems predict hazard presence, NOT impact",
        "No intelligence on which villages cut off or alternate routes"
    ], RISK_RED)
]

card_w = Inches(3.7)
card_h = Inches(2.3)
for i, (col_title, items, border_c) in enumerate(cols):
    left = Inches(0.8 + i * 4.0)
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.35), card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = border_c
    card.line.width = Pt(1.5)
    c_tf = card.text_frame
    c_tf.word_wrap = True
    c_tf.margin_left = Inches(0.2)
    c_tf.margin_top = Inches(0.15)
    
    p = c_tf.paragraphs[0]
    p.text = col_title
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = border_c
    
    for itm in items:
        p = c_tf.add_paragraph()
        p.text = f"• {itm}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_WHITE

# Visual Comparison: Current Situation vs Our Vision
comp_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.05))
comp_box.fill.solid()
comp_box.fill.fore_color.rgb = CARD_BG
comp_box.line.color.rgb = CARD_BORDER
comp_box.line.width = Pt(1)
comp_tf = comp_box.text_frame
comp_tf.word_wrap = True
comp_tf.margin_left = Inches(0.25)
comp_tf.margin_top = Inches(0.15)

p = comp_tf.paragraphs[0]
p.text = "PARADIGM SHIFT: REACTIVE TRAGEDY  ➔  PROACTIVE PREPAREDNESS"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = TEXT_ACCENT

# Current vs Vision boxes inside
curr_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(5.35), Inches(5.2), Inches(1.3))
curr_box.fill.solid()
curr_box.fill.fore_color.rgb = RGBColor(0x2A, 0x12, 0x15)
curr_box.line.color.rgb = RISK_RED
curr_box.line.width = Pt(1)
cb_tf = curr_box.text_frame
cb_tf.word_wrap = True
p = cb_tf.paragraphs[0]
p.text = "CURRENT SITUATION (REACTIVE)"
p.font.bold = True
p.font.size = Pt(11)
p.font.color.rgb = RISK_RED
p = cb_tf.add_paragraph()
p.text = "Detect Hazard  ➔  Disaster Strikes  ➔  Chaos & Isolation  ➔  Delayed Reaction"
p.font.bold = True
p.font.size = Pt(10)
p.font.color.rgb = TEXT_WHITE
p = cb_tf.add_paragraph()
p.text = "• Villages isolated for weeks | Medical emergencies stranded | Zero impact forecast"
p.font.size = Pt(8.5)
p.font.color.rgb = TEXT_MUTED

vis_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(5.35), Inches(5.4), Inches(1.3))
vis_box.fill.solid()
vis_box.fill.fore_color.rgb = RGBColor(0x0C, 0x2A, 0x24)
vis_box.line.color.rgb = RISK_GREEN
vis_box.line.width = Pt(1)
vb_tf = vis_box.text_frame
vb_tf.word_wrap = True
p = vb_tf.paragraphs[0]
p.text = "OUR VISION: LANDSIGHT AI (PROACTIVE)"
p.font.bold = True
p.font.size = Pt(11)
p.font.color.rgb = RISK_GREEN
p = vb_tf.add_paragraph()
p.text = "Predict Multi-Risk  ➔  Pre-Alert Authorities  ➔  Reroute Lifelines  ➔  Prompt Relief"
p.font.bold = True
p.font.size = Pt(10)
p.font.color.rgb = TEXT_WHITE
p = vb_tf.add_paragraph()
p.text = "• 6-12 hr pre-warning | Alternate routes mapped | Ambulances & supplies pre-deployed"
p.font.size = Pt(8.5)
p.font.color.rgb = TEXT_ACCENT

# ==========================================
# SLIDE 3: PROPOSED SOLUTION
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_header(slide3, "PROPOSED SOLUTION: Introducing LANDSIGHT AI Intelligent Ecosystem")

# Left Column: The 4-Tier End-to-End Visual Workflow
wf_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.35))
wf_card.fill.solid()
wf_card.fill.fore_color.rgb = CARD_BG
wf_card.line.color.rgb = CARD_BORDER
wf_card.line.width = Pt(1)
wf_tf = wf_card.text_frame
wf_tf.word_wrap = True
wf_tf.margin_left = Inches(0.25)
wf_tf.margin_top = Inches(0.2)

p = wf_tf.paragraphs[0]
p.text = "INTELLIGENT DATA & DECISION PIPELINE"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

steps = [
    ("1. MULTI-SOURCE INPUT FUSION", "IMD Rainfall + SMAP Soil Moisture + SRTM Slope/Elevation + Sentinel-2 Satellite Imagery + Historical GSI Landslides + Offline Citizen Reports", TEXT_ACCENT),
    ("2. AI DYNAMIC RISK ENGINE", "Machine learning ensemble (XGBoost + Random Forest + Spatial CNN) processes live environmental triggers into localized risk coefficients", TEXT_WHITE),
    ("3. GIS RISK MAPPING & SPATIAL OVERLAY", "Dynamic polygon risk heatmaps rendered via PostGIS & Mapbox/Leaflet with 30m spatial resolution across vulnerable corridors", TEXT_WHITE),
    ("4. CONNECTIVITY & IMPACT INTELLIGENCE", "PostGIS Dijkstra graph engine calculates road blockages, isolated village clusters, hospital accessibility & alternate emergency corridors", ACCENT_GOLD),
    ("5. MULTI-STAKEHOLDER ACTION OUTPUTS", "Automated SMS/FCM early warnings to citizens, NDMA/SDRF priority dashboard, and offline mobile guidance for field responders", RISK_GREEN)
]

for title, desc, col in steps:
    p = wf_tf.add_paragraph()
    p.text = f"▼  {title}"
    p.font.bold = True
    p.font.size = Pt(10.5)
    p.font.color.rgb = col
    p_sub = wf_tf.add_paragraph()
    p_sub.text = desc
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_MUTED

# Right Column: Dynamic Risk Levels & GIS Representation
risk_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.5), Inches(4.0), Inches(5.35))
risk_card.fill.solid()
risk_card.fill.fore_color.rgb = CARD_BG
risk_card.line.color.rgb = CARD_BORDER
risk_card.line.width = Pt(1)
rc_tf = risk_card.text_frame
rc_tf.word_wrap = True
rc_tf.margin_left = Inches(0.2)
rc_tf.margin_top = Inches(0.2)

p = rc_tf.paragraphs[0]
p.text = "DYNAMIC RISK CLASSIFICATION"
p.font.size = Pt(12.5)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE

risk_levels = [
    ("GREEN: LOW RISK (0 - 30%)", "Normal terrain stability. Regular automated satellite & weather polling. No disruption to road traffic.", RISK_GREEN),
    ("YELLOW: MODERATE (31 - 60%)", "Elevated saturation or light rainfall. Advisory issued to highway patrols; field officers alerted to watch slopes.", RISK_YELLOW),
    ("ORANGE: HIGH RISK (61 - 80%)", "Heavy precipitation + steep slope trigger. Pre-warning alerts to transport authorities; heavy vehicles restricted.", RISK_ORANGE),
    ("RED: CRITICAL RISK (81 - 100%)", "Imminent landslide danger. Instant Siren/SMS push, bypass corridor active, NDRF/SDRF mobilization ordered.", RISK_RED)
]

for title, desc, col in risk_levels:
    p = rc_tf.add_paragraph()
    p.text = f"■  {title}"
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = col
    p_sub = rc_tf.add_paragraph()
    p_sub.text = desc
    p_sub.font.size = Pt(8.5)
    p_sub.font.color.rgb = TEXT_WHITE

# ==========================================
# SLIDE 4: INNOVATIONS / USP
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)
add_header(slide4, "WHY LANDSIGHT AI IS DIFFERENT: Key Innovations & USPs")

# 4 Innovation Cards in 2x2 Grid
cards = [
    ("CARD 1: MULTI-SOURCE AI RISK ENGINE", 
     "Holistic Data Fusion Beyond Simple Rain Gauges",
     [
         "Synthesizes 6 distinct environmental & human streams",
         "Combines static factors (slope, geology, soil type) with dynamic triggers (24h/72h rainfall, moisture saturation)",
         "Eliminates false alarms via multi-sensor cross-validation"
     ], ACCENT_CYAN, Inches(0.8), Inches(1.5)),
    
    ("CARD 2: CONNECTIVITY IMPACT INTELLIGENCE",
     "Beyond Prediction: Answering 'What Gets Cut Off?'",
     [
         "Autonomous Road Blockage & Isolation Prediction",
         "Pinpoints isolated villages & cut-off primary health centers",
         "Auto-computes safest alternative emergency corridors (Dijkstra)",
         "Prioritizes emergency dispatch before disaster strikes"
     ], ACCENT_GOLD, Inches(6.8), Inches(1.5)),
    
    ("CARD 3: OFFLINE COMMUNITY & FIELD REPORTING",
     "Zero-Internet Ground Intelligence for Remote Valleys",
     [
         "Lightweight Flutter Mobile App with local SQLite cache",
         "Locals & officers log ground cracks, rockfalls & slope slips",
         "Geotagged photos + EXIF metadata queued offline",
         "Auto-syncs instantly when connection is restored"
     ], RISK_GREEN, Inches(0.8), Inches(4.3)),
    
    ("CARD 4: EXPLAINABLE AI (XAI) FOR DECISION TRUST",
     "Transparent Reasoning for District Disaster Authorities",
     [
         "Demystifies black-box ML using SHAP factor attributions",
         "Concrete audit: 'CRITICAL RISK — 87%' Breakdown:",
         "• 42% Extreme 24h Rainfall (185mm)  |  • 21% Soil Moisture (92%)",
         "• 15% Slope Angle (>48°)  |  • 9% Recent Citizen Crack Reports"
     ], RISK_RED, Inches(6.8), Inches(4.3))
]

c_w = Inches(5.7)
c_h = Inches(2.65)
for title, subtitle, points, col, l, t in cards:
    c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, c_w, c_h)
    c.fill.solid()
    c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = col
    c.line.width = Pt(1.5)
    ctf = c.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.2)
    ctf.margin_top = Inches(0.12)
    
    p = ctf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = col
    
    p = ctf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(9.5)
    p.font.italic = True
    p.font.color.rgb = TEXT_WHITE
    
    for pt in points:
        p = ctf.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(8.5)
        p.font.color.rgb = TEXT_MUTED if not pt.startswith("•") else TEXT_ACCENT

# ==========================================
# SLIDE 5: SYSTEM ARCHITECTURE & TECHNOLOGY
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_header(slide5, "HOW IT WORKS: 4-Tier System Architecture & Technology Stack")

# 4 Horizontal Architecture Layers
layers = [
    ("LAYER 1: DATA INGESTION & SENSING", 
     "IMD AWS Rain Data • NASA/ISRO SMAP Soil Moisture • SRTM 30m DEM Terrain • Sentinel-2 Multispectral Satellite • GSI Bhukosh Landslide Inventory • Mobile Citizen & Field Officer Reports",
     ACCENT_CYAN),
    ("LAYER 2: DATA PROCESSING & AI ANALYTICS ENGINE",
     "FastAPI Processing Pipeline • Pandas/NumPy Feature Engineering • ML Risk Ensemble (XGBoost + Random Forest) • SHAP Explainable AI Module • 30m Pixel-Level Risk Classification",
     ACCENT_BLUE),
    ("LAYER 3: GEOSPATIAL & CONNECTIVITY IMPACT ENGINE",
     "PostgreSQL + PostGIS Spatial Database • pgRouting Network Graph Topology • Dijkstra Emergency Rerouting Engine • Village Isolation & Critical Infrastructure Vulnerability Mapper",
     ACCENT_GOLD),
    ("LAYER 4: MULTI-CHANNEL PRESENTATION & DISPATCH",
     "Authority Command Dashboard (React.js + Mapbox GL) • Citizen & Responder Mobile App (Flutter + SQLite Offline) • Firebase Cloud Messaging Push + CAP-compliant SMS Broadcast Alerts",
     RISK_GREEN)
]

for i, (l_title, l_desc, l_col) in enumerate(layers):
    l_top = Inches(1.45 + i * 1.15)
    layer_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), l_top, Inches(11.7), Inches(1.0))
    layer_box.fill.solid()
    layer_box.fill.fore_color.rgb = CARD_BG
    layer_box.line.color.rgb = l_col
    layer_box.line.width = Pt(1.5)
    ltf = layer_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0.2)
    ltf.margin_top = Inches(0.1)
    
    p = ltf.paragraphs[0]
    p.text = l_title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = l_col
    
    p = ltf.add_paragraph()
    p.text = l_desc
    p.font.size = Pt(9.2)
    p.font.color.rgb = TEXT_WHITE

# Tech Stack Clean Grid at the Bottom
tech_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.85))
tech_box.fill.solid()
tech_box.fill.fore_color.rgb = CARD_BG
tech_box.line.color.rgb = CARD_BORDER
tech_box.line.width = Pt(1)
ttf = tech_box.text_frame
ttf.word_wrap = True
ttf.margin_left = Inches(0.2)
ttf.margin_top = Inches(0.1)

p = ttf.paragraphs[0]
p.text = "TARGET TECHNOLOGY STACK (FEASIBLE & PURPOSE-BUILT):"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = TEXT_ACCENT

p = ttf.add_paragraph()
p.text = "Frontend: React.js + Mapbox GL | Mobile: Flutter (Android/iOS) + SQLite Local Cache | Backend: Python + FastAPI | ML: Scikit-learn + SHAP | Spatial DB: PostgreSQL + PostGIS | Cloud: AWS EC2 + S3 | Alerts: Firebase Cloud Messaging"
p.font.size = Pt(8.8)
p.font.color.rgb = TEXT_WHITE

# ==========================================
# SLIDE 6: IMPACT, FEASIBILITY & FUTURE
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_header(slide6, "IMPACT BEYOND PREDICTION: Feasibility & Roadmap")

# 3 Pillars on Slide 6
pillars = [
    ("EXPECTED IMPACT", [
        "6-12 Hour Warning Window: Replaces post-disaster rescue with proactive community evacuation",
        "Safeguards Key Lifelines: Prevents passenger stranding along NH-10 (Sikkim) & NH-29 (Nagaland)",
        "Zero Blindspots: Offline sync empowers 100+ isolated tribal hill hamlets",
        "Data-Driven Authority Decisions: NDMA/SDMA allocate NDRF teams to high-vulnerability choke points",
        "Economic Continuity: Cuts multi-crore highway clearance and transport standstill losses"
    ], RISK_GREEN, Inches(0.8)),
    
    ("FEASIBILITY & STUDENT MVP", [
        "Open Datasets Ready: NASA SMAP, SRTM DEM, IMD Gridded Rainfall, and GSI Bhukosh are open",
        "Modular Microservice Architecture: Allows rapid independent development across 6 team members",
        "Focused Pilot District: Prototype validated on high-risk corridor (e.g., Gangtok-Mangan or Dima Hasao)",
        "Works with Low Hardware Cost: Uses cloud compute without requiring expensive localized physical IoT sensors",
        "Strictly Realistic Claims: AI acts as decision-support platform, not an infallible oracle"
    ], ACCENT_BLUE, Inches(4.75)),
    
    ("FUTURE HORIZONS", [
        "Drone Photogrammetry: Automated edge drone inspection along detected critical slope fractures",
        "InSAR Radar Satellite: Millimeter-level ground subsidence tracking via Sentinel-1 interferometry",
        "Emergency Services Integration: Direct API integration with 112 India, BRO, and District Emergency Centers",
        "Regional Language Support: Multilingual audio alerts in Assamese, Bengali, Mizo, Khasi, Hindi, Nepali",
        "Pan-Himalayan Expansion: Extension to Uttarakhand, Himachal Pradesh, and Jammu & Kashmir"
    ], ACCENT_GOLD, Inches(8.7))
]

p_w = Inches(3.8)
p_h = Inches(4.65)
for title, points, col, left in pillars:
    card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), p_w, p_h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = col
    card.line.width = Pt(1.5)
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.2)
    ctf.margin_top = Inches(0.15)
    
    p = ctf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = col
    
    for pt in points:
        p = ctf.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(8.8)
        p.font.color.rgb = TEXT_WHITE

# Bottom Tagline Banner
tag_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.7))
tag_box.fill.solid()
tag_box.fill.fore_color.rgb = RGBColor(0x0F, 0x2A, 0x4A)
tag_box.line.color.rgb = ACCENT_CYAN
tag_box.line.width = Pt(1)
t_tf = tag_box.text_frame
t_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = t_tf.paragraphs[0]
p.text = "LANDSIGHT AI:  “Predict Risk. Protect Lives. Connect Communities.”"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = TEXT_ACCENT
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = "SIH2026_LANDSIGHT_AI_Official_6Slides.pptx"
prs.save(output_path)
print(f"Presentation successfully saved to: {output_path}")
